#!/usr/bin/env python3
"""Reproduce the quantitative results of the shared-micromobility lifecycle
data-exposure study: all figures, the editable tables, and every count,
proportion, and confidence interval reported in the article.

This module is the public reproducibility engine. It depends only on the
committed public data (``data/``, ``results/``, ``review/``) and regenerates the
figures and the in-text numbers with a single command::

    python reproduce.py

It deliberately contains no manuscript prose or cover-letter text: the
manuscript-body assembler lives in a separate, non-public script and imports
the loaders, figures, and tables defined here so that the published numbers and
the reproducible numbers can never diverge.
"""

from __future__ import annotations

import json

import csv
import re
import shutil
import textwrap
from collections import Counter
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
import numpy as np

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "output"
WORK = ROOT / "build"
FIGDIR = OUTPUT / "figures"


# ---------------------------------------------------------------------------
# References. Keyed by short label; rendered in author-date (Cambridge A)
# style and alphabetised in the reference list (see resolve_citations). Every
# reference below was checked against a public identifier (DOI/URL) during
# preparation and is recorded in output/Reference_Verification.csv.
# ---------------------------------------------------------------------------
REFS: dict[str, str] = {
    "prisma_scr": (
        "Tricco, A.C., Lillie, E., Zarin, W., O'Brien, K.K., Colquhoun, H., "
        "Levac, D., et al. (2018) PRISMA Extension for Scoping Reviews "
        "(PRISMA-ScR): checklist and explanation. Annals of Internal Medicine "
        "169(7), 467-473. https://doi.org/10.7326/M18-0850."
    ),
    "grade_indirectness": (
        "Guyatt, G.H., Oxman, A.D., Kunz, R., Woodcock, J., Brozek, J., "
        "Helfand, M., et al. (2011) GRADE guidelines: 8. Rating the quality of "
        "evidence - indirectness. Journal of Clinical Epidemiology 64(12), "
        "1303-1310. https://doi.org/10.1016/j.jclinepi.2011.04.014."
    ),
    "arksey": (
        "Arksey, H. and O'Malley, L. (2005) Scoping studies: towards a "
        "methodological framework. International Journal of Social Research "
        "Methodology 8(1), 19-32. https://doi.org/10.1080/1364557032000119616."
    ),
    "demontjoye": (
        "de Montjoye, Y.-A., Hidalgo, C.A., Verleysen, M. and Blondel, V.D. "
        "(2013) Unique in the crowd: the privacy bounds of human mobility. "
        "Scientific Reports 3, 1376. https://doi.org/10.1038/srep01376."
    ),
    "gdpr": (
        "European Parliament and Council (2016) Regulation (EU) 2016/679 "
        "(General Data Protection Regulation). Official Journal of the European "
        "Union L119, 1-88."
    ),
    "elzer": (
        "Elzer, T., Ruben, M., Classen, J. and Hollick, M. (2025) They see me "
        "scooting: a long-term real-world data analysis of shared "
        "micro-mobility services. In 2025 IEEE 10th European Symposium on "
        "Security and Privacy (EuroS&P). IEEE. "
        "https://doi.org/10.1109/EuroSP63326.2025.00014."
    ),
    "etrojans": (
        "Casagrande, M., Losiouk, E., Conti, M., Payer, M. and Antonioli, D. "
        "(2024) E-Trojans: ransomware, tracking, DoS, and data leaks on "
        "battery-powered embedded systems. arXiv:2411.17184. "
        "https://doi.org/10.48550/arXiv.2411.17184."
    ),
    "espoofer": (
        "Casagrande, M., Losiouk, E., Conti, M., Payer, M. and Antonioli, D. "
        "(2023) E-Spoofer: attacking and defending Xiaomi electric scooter "
        "ecosystem. In Proceedings of the 16th ACM Conference on Security and "
        "Privacy in Wireless and Mobile Networks (WiSec). ACM. "
        "https://doi.org/10.1145/3558482.3590176."
    ),
    "vinayaga2022": (
        "Vinayaga-Sureshkanth, N., Maiti, A., Jadliwala, M., Crager, K., He, J. "
        "and Rathore, H. (2022) An investigative study on the privacy "
        "implications of mobile e-scooter rental apps. In Proceedings of the "
        "15th ACM Conference on Security and Privacy in Wireless and Mobile "
        "Networks (WiSec). ACM. https://doi.org/10.1145/3507657.3528551."
    ),
    "hilgert": (
        "Hilgert, J.-N., Lambertz, M., Rybalka, M., Schell, R. and Vogt, R. "
        "(2021) A forensic analysis of micromobility solutions. Forensic "
        "Science International: Digital Investigation 38, 301137. "
        "https://doi.org/10.1016/j.fsidi.2021.301137."
    ),
    "isik": (
        "Isik, A.B., Dag, T. and Ozkan, K. (2023) E-scooter sharing platforms: "
        "understanding their architecture and cybersecurity threats. In 2023 "
        "IEEE 26th International Conference on Intelligent Transportation "
        "Systems (ITSC). IEEE. https://doi.org/10.1109/ITSC57777.2023.10421849."
    ),
    "vinayaga2020": (
        "Vinayaga-Sureshkanth, N., Maiti, A., Jadliwala, M., Crager, K., He, J. "
        "and Rathore, H. (2020) Security and privacy challenges in upcoming "
        "intelligent urban micromobility transportation systems. In Proceedings "
        "of the Second ACM Workshop on Automotive and Aerial Vehicle Security "
        "(AutoSec). ACM. https://doi.org/10.1145/3375706.3380559."
    ),
    "petersen": (
        "Petersen, M.L. (2019) Scoot over smart devices: the invisible costs of "
        "rental scooters. Surveillance & Society 17(1/2), 267-273. "
        "https://doi.org/10.24908/ss.v17i1/2.13112."
    ),
    "sato": (
        "Sato, K., Fukushima, N., Fujii, K. and Kitani, T. (2025) Data "
        "acquisition framework for micromobility vehicles toward driving risk "
        "prediction. IEEE Security & Privacy 23(1). "
        "https://doi.org/10.1109/MSEC.2024.3441731."
    ),
    "yilmaz2022": (
        "Yilmaz, S. and Karsligil, M.E. (2022) Analysis of location spoofing "
        "threats on e-scooter sharing. In 2022 30th Signal Processing and "
        "Communications Applications Conference (SIU). IEEE. "
        "https://doi.org/10.1109/SIU55565.2022.9864946."
    ),
    "yilmaz2023": (
        "Yilmaz, S. and Karsligil, M.E. (2023) Geo-location spoofing on "
        "e-scooters: threat analysis and prevention framework. Balkan Journal "
        "of Electrical and Computer Engineering 11(2). "
        "https://doi.org/10.17694/bajece.1231384."
    ),
    "li2020": (
        "Li, Y. and Zhang, X. (2020) Linking privacy concerns for traceable "
        "information and information privacy protective responses. In "
        "Proceedings of the 53rd Hawaii International Conference on System "
        "Sciences (HICSS). https://doi.org/10.24251/HICSS.2020.105."
    ),
    "hannemann": (
        "Hannemann, A., Buchholz, E. and Ziegler, D. (2023) Is homomorphic "
        "encryption feasible for smart mobility? In Annals of Computer Science "
        "and Information Systems, vol. 35 (FedCSIS). "
        "https://doi.org/10.15439/2023F695."
    ),
    "zhou": (
        "Zhou, Y. (2024) Data extraction in dockless bikeshare: an analysis "
        "from users' perspective. Big Data & Society 11(4). "
        "https://doi.org/10.1177/20539517241299724."
    ),
    "leaky": (
        "Marchiori, A., Losiouk, E., Conti, M. and Antonioli, D. (2025) Leaky "
        "batteries: a novel set of side-channel attacks on electric vehicles. "
        "In Computer Security - ESORICS 2025, Lecture Notes in Computer "
        "Science. Springer. https://doi.org/10.1007/978-3-032-00624-0_16."
    ),
    "iotreuse": (
        "Liu, Y., Zhou, Z., Zhang, J., et al. (2023) How IoT re-using threatens "
        "your sensitive data: exploring the user-data disposal in used IoT "
        "devices. In 2023 IEEE Symposium on Security and Privacy (SP). IEEE. "
        "https://doi.org/10.1109/SP46215.2023.10179294."
    ),
    "bms": (
        "Basic, F., Gaertner, M. and Steger, C. (2023) Secure data acquisition "
        "for battery management systems. In 2023 26th Euromicro Conference on "
        "Digital System Design (DSD). IEEE. "
        "https://doi.org/10.1109/DSD60849.2023.00082."
    ),
    "remanence": (
        "Joshi, A. and Raval, M.S. (2018) Standards and techniques to remove "
        "data remanence in cloud storage. In 2018 IEEE Punecon. IEEE. "
        "https://doi.org/10.1109/PUNECON.2018.8745370."
    ),
    "gbfs_spec": (
        "MobilityData (2024) General Bikeshare Feed Specification (GBFS), "
        "version 3.0. Available at https://gbfs.org/ (accessed 12 July 2026)."
    ),
    "gbfs_registry": (
        "MobilityData (2026) GBFS systems catalogue (systems.csv). Available at "
        "https://github.com/MobilityData/gbfs (accessed 8 July 2026)."
    ),
    "mds_privacy": (
        "Open Mobility Foundation (2021) MDS privacy guide for cities. "
        "Available at https://www.openmobilityfoundation.org/ (accessed 12 "
        "July 2026)."
    ),
    "edpb_cv": (
        "European Data Protection Board (2021) Guidelines 01/2020 on processing "
        "personal data in the context of connected vehicles and mobility "
        "related applications, version 2.0."
    ),
    "nist88": (
        "Kissel, R., Regenscheid, A., Scholl, M. and Stine, K. (2014) NIST "
        "Special Publication 800-88 Revision 1: Guidelines for Media "
        "Sanitization. National Institute of Standards and Technology. "
        "https://doi.org/10.6028/NIST.SP.800-88r1."
    ),
    "nist161": (
        "Boyens, J., Smith, A., Bartol, N., Winkler, K., Holbrook, A. and "
        "Fallon, M. (2022) NIST Special Publication 800-161 Revision 1: "
        "Cybersecurity Supply Chain Risk Management Practices for Systems and "
        "Organizations. National Institute of Standards and Technology. "
        "https://doi.org/10.6028/NIST.SP.800-161r1."
    ),
    "eu_battery": (
        "European Parliament and Council (2023) Regulation (EU) 2023/1542 "
        "concerning batteries and waste batteries. Official Journal of the "
        "European Union L191, 1-117."
    ),
}

# In-text author-date citation metadata (Cambridge A style). Each entry gives
# the in-text author string, the year, and an alphabetical sort key (first
# author surname / organisation) for ordering the reference list.
CITEMETA: dict[str, tuple[str, str, str]] = {
    "prisma_scr": ("Tricco et al.", "2018", "tricco"),
    "grade_indirectness": ("Guyatt et al.", "2011", "guyatt"),
    "arksey": ("Arksey and O'Malley", "2005", "arksey"),
    "demontjoye": ("de Montjoye et al.", "2013", "de montjoye"),
    "gdpr": ("European Parliament and Council", "2016", "european parliament and council"),
    "elzer": ("Elzer et al.", "2025", "elzer"),
    "etrojans": ("Casagrande et al.", "2024", "casagrande"),
    "espoofer": ("Casagrande et al.", "2023", "casagrande"),
    "vinayaga2022": ("Vinayaga-Sureshkanth et al.", "2022", "vinayaga-sureshkanth"),
    "hilgert": ("Hilgert et al.", "2021", "hilgert"),
    "isik": ("Isik et al.", "2023", "isik"),
    "vinayaga2020": ("Vinayaga-Sureshkanth et al.", "2020", "vinayaga-sureshkanth"),
    "petersen": ("Petersen", "2019", "petersen"),
    "sato": ("Sato et al.", "2025", "sato"),
    "yilmaz2022": ("Yilmaz and Karsligil", "2022", "yilmaz"),
    "yilmaz2023": ("Yilmaz and Karsligil", "2023", "yilmaz"),
    "li2020": ("Li and Zhang", "2020", "li"),
    "hannemann": ("Hannemann et al.", "2023", "hannemann"),
    "zhou": ("Zhou", "2024", "zhou"),
    "leaky": ("Marchiori et al.", "2025", "marchiori"),
    "iotreuse": ("Liu et al.", "2023", "liu"),
    "bms": ("Basic et al.", "2023", "basic"),
    "remanence": ("Joshi and Raval", "2018", "joshi"),
    "gbfs_spec": ("MobilityData", "2024", "mobilitydata"),
    "gbfs_registry": ("MobilityData", "2026", "mobilitydata"),
    "mds_privacy": ("Open Mobility Foundation", "2021", "open mobility foundation"),
    "edpb_cv": ("European Data Protection Board", "2021", "european data protection board"),
    "nist88": ("Kissel et al.", "2014", "kissel"),
    "nist161": ("Boyens et al.", "2022", "boyens"),
    "eu_battery": ("European Parliament and Council", "2023", "european parliament and council"),
}

assert set(CITEMETA) == set(REFS), (
    "CITEMETA and REFS must cover the same labels: "
    f"{set(CITEMETA) ^ set(REFS)}"
)


def intext(label: str) -> str:
    authors, year, _ = CITEMETA[label]
    return f"{authors} {year}"


def sortkey(label: str) -> tuple[str, str]:
    _, year, srt = CITEMETA[label]
    return (srt, year)

# ---------------------------------------------------------------------------
# Data access helpers
# ---------------------------------------------------------------------------
def load_gbfs_summary() -> dict[tuple[str, str], dict]:
    table: dict[tuple[str, str], dict] = {}
    with (ROOT / "results" / "gbfs_summary.csv").open(encoding="utf-8") as fh:
        for row in csv.DictReader(fh):
            table[(row["stratum"], row["metric"])] = row
    return table


def load_document_audit() -> tuple[list[str], dict[str, dict[str, str]], int, int]:
    rows = list(csv.DictReader((ROOT / "data" / "document_audit.csv").open(encoding="utf-8")))
    domains: list[str] = []
    matrix: dict[str, dict[str, str]] = {}
    operators_order: list[str] = []
    for r in rows:
        op = r["operator_name"]
        if op not in matrix:
            matrix[op] = {}
            operators_order.append(op)
        matrix[op][r["domain"]] = r["status"]
        if r["domain"] not in domains:
            domains.append(r["domain"])
    coded = sum(1 for op in operators_order
                if any(v != "unavailable" for v in matrix[op].values()))
    return domains, matrix, len(operators_order), coded


def load_extraction() -> list[dict]:
    return list(csv.DictReader((ROOT / "review" / "evidence_extraction.csv").open(encoding="utf-8")))


GBFS = load_gbfs_summary()
DOMAINS, DOCMATRIX, N_OPERATORS, N_CODED = load_document_audit()
EXTRACTION = load_extraction()


def g(stratum: str, metric: str, field: str) -> str:
    return GBFS[(stratum, metric)][field]


def pct(stratum: str, metric: str) -> str:
    row = GBFS[(stratum, metric)]
    return f"{float(row['percent']):.1f}%"


def ci(stratum: str, metric: str) -> str:
    row = GBFS[(stratum, metric)]
    return f"{float(row['wilson_95_low_percent']):.1f}-{float(row['wilson_95_high_percent']):.1f}%"


def domain_count(domain: str, status: str) -> int:
    return sum(1 for op, cells in DOCMATRIX.items() if cells.get(domain) == status)


def load_screening_stats() -> dict[str, int]:
    """Derive every PRISMA-flow count directly from the committed screening
    sheet so that no flow number is hard-coded in the manuscript."""
    rows = list(csv.DictReader(
        (ROOT / "review" / "screening.csv").open(encoding="utf-8")))
    ta = Counter(r["title_abstract_decision"] for r in rows)
    excl = Counter(r["full_text_exclusion_reason"] for r in rows
                   if r["full_text_decision"] == "exclude")
    included = [r for r in rows if r["full_text_decision"] == "include"]
    dist = Counter(r["evidence_distance"] for r in included)
    second_pass = [r for r in rows if r["second_pass_decision"].strip()]
    agree = sum(1 for r in second_pass
                if r["second_pass_decision"].strip() ==
                r["title_abstract_decision"].strip())
    n_no_data_path = excl["F1_NO_RELEVANT_DATA_PATH"]
    n_not_transferable = excl["F2_MECHANISM_NOT_TRANSFERABLE"]
    return {
        "identified": len(rows),
        "ta_excluded": ta["exclude"],
        "sought": ta["include"] + ta["uncertain"],
        "not_retrieved": excl["F4_NOT_RETRIEVED"],
        "no_data_path": n_no_data_path,
        "not_transferable": n_not_transferable,
        "excluded_fulltext": n_no_data_path + n_not_transferable,
        "included": len(included),
        "with_doi": sum(1 for r in rows if r["doi"].strip()),
        "d4": dist["D4"], "d3": dist["D3"], "d2": dist["D2"],
        "resample_n": len(second_pass),
        "resample_agreement_pct":
            round(100 * agree / len(second_pass)) if second_pass else 0,
    }


def _summary_int(stratum: str, metric: str, field: str) -> int:
    return int(g(stratum, metric, field))


SCR = load_screening_stats()

# Audit totals, all read from the reproducible GBFS summary rather than typed.
N_REGISTRY = _summary_int(
    "all_registry_entries", "auto_discovery_reachable", "denominator")
N_REACHABLE = _summary_int(
    "all_registry_entries", "auto_discovery_reachable", "count")
N_MOTOR_FEEDS = _summary_int(
    "declared_motorized_micromobility_feeds", "has_vehicle_id", "denominator")
N_OP_DOMAINS = _summary_int(
    "declared_motorized_micromobility_operator_domains_any",
    "has_vehicle_id", "denominator")


def fmt(n: int) -> str:
    """Render an integer with thousands separators (e.g. 2169 -> '2,169')."""
    return f"{n:,}"

# ---------------------------------------------------------------------------
# Figures (English text only, per manuscript language rule)
# ---------------------------------------------------------------------------
FIG_CAPTIONS = {
    1: ("Figure 1. PRISMA-ScR flow of record identification, title/abstract "
        "screening, full-text eligibility assessment, and inclusion. Counts are "
        "reproduced from the committed screening sheet."),
    2: (f"Figure 2. Evidence map of the {SCR['included']} included direct studies "
        "by evidence distance (D4 direct empirical, D3 direct documentary, D2 "
        "near-domain empirical) and lifecycle stage addressed. A study that "
        "addresses more than one stage is counted in each relevant stage, so row "
        "totals may exceed the number of studies at that distance."),
    3: (f"Figure 3. Prevalence of disclosed fields in {fmt(N_MOTOR_FEEDS)} "
        "successfully retrieved, non-empty public vehicle-status feeds that "
        "declared motorized micromobility. Bars show the point estimate; "
        "whiskers show the 95% Wilson confidence interval. Field presence is a "
        "disclosure signal, not evidence of a privacy harm."),
    4: (f"Figure 4. Disclosure-domain coding of public operator privacy notices "
        f"across {len(DOMAINS)} domains (n = {N_CODED} coded operators; "
        f"{N_OPERATORS - N_CODED} further operators were unavailable for "
        "reproducible retrieval). Cells show explicit, partial, or not-found "
        "coding; not-found denotes silence in the document, not evidence that a "
        "practice is absent."),
    5: ("Figure 5. Lifecycle data-exposure model linking each stage to the "
        "evidence distance of its supporting sources and to a proposed, as-yet "
        "unvalidated control. Solid boxes denote stages with direct (D3-D4) "
        "evidence; dashed boxes denote stages supported only by near-domain "
        "(D2) evidence."),
}

PALETTE = {
    "explicit": "#2c7fb8",
    "partial": "#7fcdbb",
    "not_found": "#edf8b1",
    "unavailable": "#d9d9d9",
    "bar": "#2c7fb8",
    "D4": "#08519c",
    "D3": "#3182bd",
    "D2": "#9ecae1",
}


def fig1_prisma(ax) -> None:
    ax.axis("off")
    boxes = [
        (0.5, 0.93, f"Records identified from\nbibliographic sources (n = {fmt(SCR['identified'])})"),
        (0.5, 0.75, f"Records screened on title/abstract\n(n = {fmt(SCR['identified'])})"),
        (0.5, 0.55, f"Records sought for full-text\nassessment (n = {fmt(SCR['sought'])})"),
        (0.5, 0.33, f"Full-text records assessed\nfor eligibility (n = {fmt(SCR['sought'])})"),
        (0.5, 0.11, f"Direct studies included in\nsynthesis (n = {fmt(SCR['included'])})"),
    ]
    for x, y, text in boxes:
        ax.add_patch(FancyBboxPatch((x - 0.22, y - 0.055), 0.44, 0.11,
                     boxstyle="round,pad=0.01", fc="#deebf7", ec="#08519c", lw=1.2,
                     transform=ax.transAxes))
        ax.text(x, y, text, ha="center", va="center", fontsize=8.5,
                transform=ax.transAxes)
    excl = [
        (0.80, 0.75, f"Excluded on title/abstract\n(n = {fmt(SCR['ta_excluded'])})"),
        (0.80, 0.44, f"Full text not retrieved\n(n = {fmt(SCR['not_retrieved'])})"),
        (0.80, 0.22, f"Excluded at full text (n = {fmt(SCR['excluded_fulltext'])}):\n"
                     f"no relevant data path (n = {fmt(SCR['no_data_path'])});\n"
                     f"mechanism not transferable (n = {fmt(SCR['not_transferable'])})"),
    ]
    for x, y, text in excl:
        ax.add_patch(FancyBboxPatch((x - 0.18, y - 0.06), 0.36, 0.12,
                     boxstyle="round,pad=0.01", fc="#f7f7f7", ec="#999999", lw=1.0,
                     transform=ax.transAxes))
        ax.text(x, y, text, ha="center", va="center", fontsize=7.5,
                transform=ax.transAxes)
    ys = [0.875, 0.695, 0.475, 0.275]
    for y0 in ys:
        ax.add_patch(FancyArrowPatch((0.5, y0), (0.5, y0 - 0.05),
                     arrowstyle="-|>", mutation_scale=12, color="#08519c",
                     transform=ax.transAxes))
    for y0, yx in [(0.75, 0.75), (0.55, 0.44), (0.33, 0.22)]:
        ax.add_patch(FancyArrowPatch((0.72, y0), (0.62, y0),
                     arrowstyle="-|>", mutation_scale=10, color="#999999",
                     transform=ax.transAxes))


def fig2_evidence_map(ax) -> None:
    stages = ["Operation", "Maintenance/\nservice", "Recall/return",
              "Second-life/\ndisposal"]
    distances = ["D2", "D3", "D4"]
    # map each study to (distance, stage) using lifecycle_stage in extraction
    stage_key = {
        "operation": 0, "data sharing": 0,
        "maintenance": 1, "maintenance and service": 1,
        "recall and return": 2, "recall": 2, "return": 2,
        "second-life and disposal": 3, "second-life": 3, "disposal": 3,
        "recycling": 3, "manufacturing": 3,
    }
    counts = np.zeros((len(distances), len(stages)), dtype=int)
    for e in EXTRACTION:
        d = distances.index(e["evidence_distance"]) if e["evidence_distance"] in distances else None
        if d is None:
            continue
        seen: set[int] = set()
        for token in e["lifecycle_stage"].split(";"):
            s = stage_key.get(token.strip().lower())
            if s is not None:
                seen.add(s)
        for s in seen:
            counts[d, s] += 1
    im = ax.imshow(counts, cmap="Blues", aspect="auto", vmin=0, vmax=max(counts.max(), 1))
    ax.set_xticks(range(len(stages)))
    ax.set_xticklabels(stages, fontsize=8)
    ax.set_yticks(range(len(distances)))
    ax.set_yticklabels(["D2 near-domain", "D3 documentary", "D4 empirical"], fontsize=8)
    for i in range(len(distances)):
        for j in range(len(stages)):
            if counts[i, j]:
                ax.text(j, i, str(counts[i, j]), ha="center", va="center",
                        color="black", fontsize=9)
    ax.set_xlabel("Lifecycle stage addressed", fontsize=8.5)
    cb = plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cb.set_label("Number of included studies", fontsize=8)


def fig3_prevalence(ax) -> None:
    metrics = [
        ("has_vehicle_id", "Vehicle\nidentifier"),
        ("has_location_fields", "Latitude/\nlongitude"),
        ("has_range", "Current\nrange"),
        ("has_deep_link", "Rental\ndeep link"),
        ("has_last_reported", "Last-reported\ntimestamp"),
        ("has_battery_percent", "Battery/fuel\npercentage"),
    ]
    stratum = "declared_motorized_micromobility_feeds"
    labels, vals, lo, hi = [], [], [], []
    for metric, lab in metrics:
        row = GBFS[(stratum, metric)]
        labels.append(lab)
        vals.append(float(row["percent"]))
        lo.append(float(row["percent"]) - float(row["wilson_95_low_percent"]))
        hi.append(float(row["wilson_95_high_percent"]) - float(row["percent"]))
    x = np.arange(len(labels))
    ax.bar(x, vals, color=PALETTE["bar"], width=0.62,
           yerr=[lo, hi], capsize=3, ecolor="#333333")
    for xi, v, h in zip(x, vals, hi):
        ax.text(xi, v + h + 1.6, f"{v:.1f}", ha="center", fontsize=8)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=8)
    ax.set_ylabel("Feeds disclosing the field (%)", fontsize=9)
    ax.set_ylim(0, 108)
    ax.set_title(
        f"Disclosed fields in {fmt(N_MOTOR_FEEDS)} motorized micromobility vehicle feeds",
        fontsize=9.5)
    ax.spines[["top", "right"]].set_visible(False)


DOMAIN_LABELS = {
    "location_data": "Location",
    "trip_time_data": "Trip/time",
    "vehicle_identifier": "Vehicle ID",
    "battery_or_diagnostic_data": "Battery/diagnostic",
    "maintenance_or_repair_data": "Maintenance/repair",
    "account_payment_device_data": "Account/payment/device",
    "analytics_or_profiling": "Analytics/profiling",
    "retention": "Retention",
    "processors_or_contractors": "Processors/contractors",
    "international_transfer": "Intl transfer",
    "user_rights": "Data-subject rights",
    "incident_contact": "Incident contact",
    "vulnerability_disclosure": "Vulnerability disclosure",
    "return_recycling_disposal": "Return/recycling/disposal",
}
STATUS_VAL = {"explicit": 3, "partial": 2, "not_found": 1, "unavailable": 0}


def fig4_heatmap(ax) -> None:
    operators = [op for op in DOCMATRIX
                 if any(v != "unavailable" for v in DOCMATRIX[op].values())]
    grid = np.array([[STATUS_VAL[DOCMATRIX[op][d]] for d in DOMAINS] for op in operators])
    from matplotlib.colors import ListedColormap
    cmap = ListedColormap([PALETTE["unavailable"], PALETTE["not_found"],
                           PALETTE["partial"], PALETTE["explicit"]])
    ax.imshow(grid, cmap=cmap, aspect="auto", vmin=0, vmax=3)
    ax.set_xticks(range(len(DOMAINS)))
    ax.set_xticklabels([DOMAIN_LABELS[d] for d in DOMAINS], rotation=45,
                       ha="right", fontsize=7)
    ax.set_yticks(range(len(operators)))
    ax.set_yticklabels(operators, fontsize=7)
    ax.set_xticks(np.arange(-0.5, len(DOMAINS), 1), minor=True)
    ax.set_yticks(np.arange(-0.5, len(operators), 1), minor=True)
    ax.grid(which="minor", color="white", linewidth=1.0)
    ax.tick_params(which="minor", length=0)
    handles = [plt.Rectangle((0, 0), 1, 1, fc=PALETTE[k]) for k in
               ["explicit", "partial", "not_found"]]
    ax.legend(handles, ["Explicit", "Partial", "Not found"],
              loc="upper left", bbox_to_anchor=(1.01, 1.0), fontsize=7,
              frameon=False)


LIFECYCLE = [
    ("Procurement", "D3", "Contract, provenance,\nand access terms"),
    ("Deployment", "D3", "Feed configuration\nand field minimization"),
    ("Operation", "D4", "Local processing;\naccess and retention limits"),
    ("Maintenance/service", "D3", "Controlled contractor\naccess and logging"),
    ("Recall/return", "D2", "Chain-of-custody\nfor returned units"),
    ("Second-life/disposal", "D2", "Verified media\nsanitization"),
]


def fig5_lifecycle(ax) -> None:
    ax.axis("off")
    n = len(LIFECYCLE)
    xs = np.linspace(0.08, 0.92, n)
    for i, ((stage, dist, control), x) in enumerate(zip(LIFECYCLE, xs)):
        direct = dist in ("D3", "D4")
        style = "round,pad=0.01" if direct else "round,pad=0.01"
        ax.add_patch(FancyBboxPatch((x - 0.065, 0.62), 0.13, 0.16,
                     boxstyle=style, fc=PALETTE.get(dist, "#deebf7"),
                     ec="#08306b", lw=1.4,
                     linestyle="-" if direct else "--", transform=ax.transAxes))
        ax.text(x, 0.70, f"{stage}\n[{dist}]", ha="center", va="center",
                fontsize=7.2, transform=ax.transAxes, color="white" if dist == "D4" else "black")
        ax.add_patch(FancyBboxPatch((x - 0.075, 0.30), 0.15, 0.18,
                     boxstyle="round,pad=0.01", fc="#f7fbff", ec="#999999",
                     lw=1.0, transform=ax.transAxes))
        ax.text(x, 0.39, control, ha="center", va="center", fontsize=6.6,
                transform=ax.transAxes)
        ax.add_patch(FancyArrowPatch((x, 0.62), (x, 0.485), arrowstyle="-|>",
                     mutation_scale=9, color="#999999", transform=ax.transAxes))
        if i < n - 1:
            ax.add_patch(FancyArrowPatch((x + 0.065, 0.70), (xs[i + 1] - 0.065, 0.70),
                         arrowstyle="-|>", mutation_scale=10, color="#08306b",
                         transform=ax.transAxes))
    ax.text(0.5, 0.90, "Lifecycle stage (with dominant evidence distance)",
            ha="center", fontsize=8.5, transform=ax.transAxes, weight="bold")
    ax.text(0.5, 0.19, "Proposed control (effectiveness not yet validated)",
            ha="center", fontsize=8.5, transform=ax.transAxes, weight="bold")


FIGFUNCS = {1: fig1_prisma, 2: fig2_evidence_map, 3: fig3_prevalence,
            4: fig4_heatmap, 5: fig5_lifecycle}
FIGSIZE = {1: (7.0, 6.0), 2: (7.0, 4.2), 3: (7.2, 4.4), 4: (8.4, 5.2), 5: (9.6, 4.6)}

def build_figures() -> dict[int, dict[str, Path]]:
    FIGDIR.mkdir(parents=True, exist_ok=True)
    paths: dict[int, dict[str, Path]] = {}
    for num, func in FIGFUNCS.items():
        fig, ax = plt.subplots(figsize=FIGSIZE[num])
        func(ax)
        fig.tight_layout()
        png = FIGDIR / f"Figure{num}.png"
        tiff = FIGDIR / f"Figure{num}.tiff"
        pdf = FIGDIR / f"Figure{num}.pdf"
        fig.savefig(png, dpi=600, bbox_inches="tight")
        fig.savefig(tiff, dpi=600, bbox_inches="tight", pil_kwargs={"compression": "tiff_lzw"})
        fig.savefig(pdf, bbox_inches="tight")
        plt.close(fig)
        paths[num] = {"png": png, "tiff": tiff, "pdf": pdf}
    return paths


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------
def build_tables() -> dict[int, dict]:
    t1 = {
        "title": "Table 1. Eligibility criteria and data sources for the three "
                 "study components.",
        "headers": ["Component", "Unit of analysis", "Inclusion criteria",
                    "Primary source"],
        "rows": [
            ["Scoping review (WP1)", "Study",
             "Peer-reviewed or archival works reporting a data path in or "
             "transferable to shared micromobility",
             f"Bibliographic search of {fmt(SCR['identified'])} records "
             "[[prisma_scr;arksey]]"],
            ["GBFS field audit (WP2)", "System/feed",
             "Public GBFS systems in the frozen registry with a reachable "
             "auto-discovery endpoint",
             "GBFS systems catalogue and live feeds [[gbfs_registry;gbfs_spec]]"],
            ["Disclosure audit (WP3)", "Operator document",
             "Publicly reachable operator privacy notice for an audited GBFS "
             "operator",
             "Operator privacy notices; governance guidance [[mds_privacy;edpb_cv]]"],
        ],
    }
    dist_label = {"D4": "D4 empirical", "D3": "D3 documentary", "D2": "D2 near-domain"}
    ref_by_id = {
        "e641ae6185ec": "espoofer", "4fe184aeb5b3": "etrojans",
        "3b9d67a99101": "elzer", "f7ef2527bcc7": "hilgert",
        "4bf4bae9ab93": "vinayaga2022", "9f85e4350278": "hannemann",
        "bdab08ae526c": "isik", "2856524cd67f": "li2020",
        "37630ff02a7b": "petersen", "b26b8d65d468": "sato",
        "ff74fa3c40ed": "vinayaga2020", "5bf50e1a9f5c": "yilmaz2022",
        "05e20798d6c6": "yilmaz2023", "62abc8d8d637": "zhou",
        "062854732ff3": "bms", "57ec900852e5": "remanence",
        "dab97e54aa73": "iotreuse", "882e5dc8fae9": "leaky",
    }
    order = {"D4": 0, "D3": 1, "D2": 2}
    ext_sorted = sorted(EXTRACTION, key=lambda e: (order[e["evidence_distance"]],
                                                   e["citation"]))
    t2_rows = []
    for e in ext_sorted:
        label = ref_by_id[e["record_id"]]
        short = e["citation"].split(".")[0]
        t2_rows.append([
            f"{short} [[year:{label}]]",
            dist_label[e["evidence_distance"]],
            e["device_or_service"],
            textwrap.shorten(e["key_result"], width=180, placeholder="..."),
        ])
    t2 = {
        "title": f"Table 2. Included direct studies (n = {len(t2_rows)}), ordered "
                 "by evidence distance. Key results are as reported by the "
                 "source; this review did not reproduce the underlying "
                 "experiments.",
        "headers": ["Study", "Evidence distance", "Device or service",
                    "Reported key result"],
        "rows": t2_rows,
    }
    def row(stratum, metric, label):
        return [label, g(stratum, metric, "count"), g(stratum, metric, "denominator"),
                pct(stratum, metric), ci(stratum, metric)]
    t3 = {
        "title": "Table 3. GBFS audit sampling and outcomes with 95% Wilson "
                 "confidence intervals. Denominators are stated explicitly; "
                 "unavailable or empty feeds are separated from absent fields.",
        "headers": ["Stratum / outcome", "n", "Denominator", "Percent", "95% CI"],
        "rows": [
            row("all_registry_entries", "auto_discovery_reachable",
                "Reachable auto-discovery endpoint"),
            row("reachable_registry_entries", "vehicle_feed_declared",
                "Vehicle-status feed declared"),
            row("successful_vehicle_feeds", "vehicle_feed_nonempty",
                "Retrieved feed with >=1 vehicle"),
            row("declared_motorized_micromobility_feeds", "has_vehicle_id",
                "Vehicle identifier (motorized feeds)"),
            row("declared_motorized_micromobility_feeds", "has_location_fields",
                "Latitude/longitude (motorized feeds)"),
            row("declared_motorized_micromobility_feeds", "has_last_reported",
                "Last-reported timestamp (motorized feeds)"),
            row("declared_motorized_micromobility_feeds", "has_battery_percent",
                "Battery/fuel percentage (motorized feeds)"),
        ],
    }
    domain_row = lambda d: [
        DOMAIN_LABELS[d], str(domain_count(d, "explicit")),
        str(domain_count(d, "partial")), str(domain_count(d, "not_found")),
    ]
    t4 = {
        "title": f"Table 4. Disclosure-domain coding of public operator privacy "
                 f"notices (n = {N_CODED} coded operators). Counts of explicit, "
                 f"partial, and not-found codings per domain; not-found denotes "
                 f"document silence, not evidence that a practice is absent.",
        "headers": ["Disclosure domain", "Explicit", "Partial", "Not found"],
        "rows": [domain_row(d) for d in DOMAINS],
    }
    t5 = {
        "title": "Table 5. Evidence-to-control traceability across the lifecycle "
                 "model. Effectiveness evidence is stated conservatively; the "
                 "proposed controls have not been empirically validated in this "
                 "study.",
        "headers": ["Lifecycle stage", "Exposure path", "Supporting evidence "
                    "(distance)", "Proposed control", "Effectiveness evidence"],
        "rows": [
            ["Procurement", "Contract and provenance terms set who can access "
             "backend data", "Governance guidance (N) [[nist161;edpb_cv]]",
             "Provenance and access clauses; data-processing terms",
             "Not validated; governance-based"],
            ["Deployment", "Feed configuration determines which fields are "
             "published", "GBFS field presence (D4) [[elzer;gbfs_spec]]",
             "Field minimization; coarse or delayed positions",
             "Not validated; design-based"],
            ["Operation", "Identifiers, positions, and timestamps enable "
             "tracking and re-identification",
             "Empirical audits (D4) [[elzer;vinayaga2022]]; mobility "
             "re-identification (D2) [[demontjoye]]",
             "Local processing; access and retention limits",
             "Partly supported by attack studies; controls not validated"],
            ["Maintenance/service", "Contractor access to devices and records",
             "Forensic and platform studies (D3-D4) [[hilgert;isik]]",
             "Controlled contractor access; audit logging",
             "Not validated"],
            ["Recall/return", "Returned units retain stored data",
             "IoT reuse study (D2) [[iotreuse]]",
             "Chain-of-custody for returned units",
             "Not validated; analogy-based"],
            ["Second-life/disposal", "Residual data on storage and batteries",
             "Data remanence and battery side-channel (D2) "
             "[[remanence;leaky;bms]]; sanitization guidance (N) [[nist88]]",
             "Verified media sanitization before resale or recycling",
             "Not validated; standards-based"],
        ],
    }
    return {1: t1, 2: t2, 3: t3, 4: t4, 5: t5}


TABLES = build_tables()

def iter_citation_texts(blocks: list[tuple], tables: dict[int, dict]):
    """Yield text fragments in document order for numbering."""
    for kind, payload in blocks:
        if kind == "p":
            yield payload
        elif kind == "table":
            spec = tables[payload]
            for r in spec["rows"]:
                for cell in r:
                    yield cell


CITE_RX = re.compile(r"\[\[([^\]]+)\]\]")


def _parse_label(raw: str) -> tuple[str, bool]:
    """Return (label, year_only). A ``year:`` prefix requests a year-only
    (narrative) citation, e.g. "Arksey and O'Malley (2005)"."""
    raw = raw.strip()
    if raw.startswith("year:"):
        return raw[len("year:"):].strip(), True
    return raw, False


def resolve_citations(blocks: list[tuple], tables: dict[int, dict]):
    """Author-date (Cambridge A) citation resolver.

    Internal ``[[label]]`` tokens are replaced with "(Author Year)" strings.
    Multiple labels in one token are sorted alphabetically by first author and
    chronologically for the same author. The reference list is returned in
    alphabetical order (not first-appearance order)."""
    order: list[str] = []

    def register(fragment: str) -> None:
        for m in CITE_RX.finditer(fragment):
            for raw in m.group(1).split(";"):
                label, _ = _parse_label(raw)
                if label not in order:
                    order.append(label)

    for frag in iter_citation_texts(blocks, tables):
        register(frag)

    unknown = [l for l in order if l not in REFS]
    if unknown:
        raise RuntimeError(f"Unknown citation labels: {unknown}")

    ref_list = sorted(set(order), key=lambda l: (sortkey(l), l))

    def repl(fragment: str) -> str:
        def _sub(m: re.Match) -> str:
            parsed = [_parse_label(p) for p in m.group(1).split(";")]
            year_only = all(yo for _, yo in parsed)
            labels = sorted((lbl for lbl, _ in parsed), key=lambda l: (sortkey(l), l))
            if year_only:
                inner = "; ".join(CITEMETA[l][1] for l in labels)
            else:
                inner = "; ".join(intext(l) for l in labels)
            return f"({inner})"
        return CITE_RX.sub(_sub, fragment)

    return order, ref_list, repl

# ---------------------------------------------------------------------------
# DOCX helpers (retained from the original generator)
# ---------------------------------------------------------------------------
def reset_dirs() -> None:
    for path in (OUTPUT, WORK):
        if path.exists():
            shutil.rmtree(path)
        path.mkdir(parents=True)


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_cell_margins(cell, top=80, start=80, bottom=80, end=80) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    tc_mar = tc_pr.first_child_found_in("w:tcMar")
    if tc_mar is None:
        tc_mar = OxmlElement("w:tcMar")
        tc_pr.append(tc_mar)
    for margin, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = tc_mar.find(qn(f"w:{margin}"))
        if node is None:
            node = OxmlElement(f"w:{margin}")
            tc_mar.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")


def configure_document(doc: Document) -> None:
    section = doc.sections[0]
    for attr in ("top_margin", "bottom_margin", "left_margin", "right_margin"):
        setattr(section, attr, Inches(1))
    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Times New Roman"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
    normal.font.size = Pt(12)
    normal.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    normal.paragraph_format.space_after = Pt(6)
    for style_name, size in (("Title", 16), ("Heading 1", 14), ("Heading 2", 12)):
        style = styles[style_name]
        style.font.name = "Times New Roman"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
        style.font.size = Pt(size)
        style.font.color.rgb = RGBColor(0, 0, 0)
    styles["Heading 1"].font.bold = True
    styles["Heading 2"].font.bold = True

def add_table(doc: Document, spec: dict, repl) -> None:
    caption = doc.add_paragraph()
    caption.paragraph_format.space_before = Pt(14)
    caption.paragraph_format.space_after = Pt(6)
    caption.add_run(repl(spec["title"])).bold = True
    table = doc.add_table(rows=1, cols=len(spec["headers"]))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    header = table.rows[0].cells
    for index, value in enumerate(spec["headers"]):
        header[index].text = value
        set_cell_shading(header[index], "D9EAF7")
        header[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        for run in header[index].paragraphs[0].runs:
            run.bold = True
            run.font.name = "Arial"
            run.font.size = Pt(8)
        set_cell_margins(header[index])
    for row_values in spec["rows"]:
        cells = table.add_row().cells
        for index, value in enumerate(row_values):
            cells[index].text = repl(value)
            cells[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.TOP
            set_cell_margins(cells[index])
            for paragraph in cells[index].paragraphs:
                paragraph.paragraph_format.space_after = Pt(0)
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(8)
    doc.add_paragraph()

def build_tables_docx(tables, repl) -> Path:
    doc = Document()
    configure_document(doc)
    h = doc.add_paragraph(style="Title")
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    h.add_run("Editable tables").bold = True
    for i in range(1, 6):
        add_table(doc, tables[i], repl)
        if i < 5:
            doc.add_section(WD_SECTION.NEW_PAGE)
    path = OUTPUT / "Tables_DataPolicy_editable.docx"
    doc.save(path)
    return path

def build_figures_pptx(figpaths, repl) -> Path:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    for num in range(1, 6):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2),
                                             PptxInches(12.3), PptxInches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Figure {num}"
        p.font.size = PptxPt(22)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        from PIL import Image as PILImage
        with PILImage.open(figpaths[num]["png"]) as im:
            w, h = im.size
        avail_w, avail_h = 12.3, 5.3
        ratio = min(avail_w / (w / 600.0), avail_h / (h / 600.0))
        disp_w = (w / 600.0) * ratio
        disp_h = (h / 600.0) * ratio
        left = (13.333 - disp_w) / 2
        slide.shapes.add_picture(str(figpaths[num]["png"]), PptxInches(left),
                                 PptxInches(1.0), PptxInches(disp_w), PptxInches(disp_h))
        cap_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(6.5),
                                          PptxInches(12.1), PptxInches(0.9))
        cf = cap_box.text_frame
        cf.word_wrap = True
        cp = cf.paragraphs[0]
        cp.text = repl(FIG_CAPTIONS[num])
        cp.font.size = PptxPt(11)
    path = OUTPUT / "Figures_DataPolicy_editable.pptx"
    prs.save(path)
    return path


# ---------------------------------------------------------------------------
# Reproducibility entry point
# ---------------------------------------------------------------------------
def dump_reproducibility_values() -> Path:
    """Write every count, proportion, and confidence interval reported in the
    article into a single machine-readable file, derived only from the committed
    public data. A reader can verify the in-text numbers from this file without
    any manuscript-generation code."""
    gbfs_rows = []
    for (stratum, metric), row in GBFS.items():
        gbfs_rows.append({
            "stratum": stratum,
            "metric": metric,
            "count": int(row["count"]),
            "denominator": int(row["denominator"]),
            "percent": float(row["percent"]),
            "wilson_95_low_percent": float(row["wilson_95_low_percent"]),
            "wilson_95_high_percent": float(row["wilson_95_high_percent"]),
        })
    operator_rows = []
    op_csv = ROOT / "results" / "gbfs_operator_summary.csv"
    if op_csv.exists():
        with op_csv.open(encoding="utf-8") as fh:
            operator_rows = list(csv.DictReader(fh))
    document_audit = {
        d: {status: domain_count(d, status)
            for status in ("explicit", "partial", "not_found", "unavailable")}
        for d in DOMAINS
    }
    payload = {
        "screening_prisma_flow": SCR,
        "gbfs_audit_totals": {
            "registry_entries": N_REGISTRY,
            "auto_discovery_reachable": N_REACHABLE,
            "motorized_micromobility_feeds": N_MOTOR_FEEDS,
            "eligible_operator_domains": N_OP_DOMAINS,
        },
        "gbfs_summary": gbfs_rows,
        "gbfs_operator_summary": operator_rows,
        "document_audit": {
            "n_operators": N_OPERATORS,
            "n_coded": N_CODED,
            "domain_counts": document_audit,
        },
    }
    path = OUTPUT / "reproducibility_values.json"
    path.write_text(json.dumps(payload, indent=2, ensure_ascii=False) + "\n",
                    encoding="utf-8")
    return path


def main() -> None:
    reset_dirs()
    figpaths = build_figures()
    _, _, repl = resolve_citations([], TABLES)
    build_tables_docx(TABLES, repl)
    build_figures_pptx(figpaths, repl)
    values_path = dump_reproducibility_values()
    print("Reproduction complete.")
    print(f"  figures: {len(figpaths)} (PNG/TIFF/PDF at 600 dpi)")
    print("  editable tables: output/Tables_DataPolicy_editable.docx")
    print("  editable figures: output/Figures_DataPolicy_editable.pptx")
    print(f"  in-text values: {values_path.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
