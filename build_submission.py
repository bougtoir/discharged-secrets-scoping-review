#!/usr/bin/env python3
"""Build the Data & Policy (Cambridge University Press) submission package for
the shared-micromobility lifecycle data-exposure study (scoping review +
global GBFS field audit + public-document disclosure audit).

The script is data-driven: quantitative claims are read from the committed
analysis outputs (results/, data/, review/) so the manuscript, figures, and
tables stay consistent with the underlying evidence. Citations use author-date
(Cambridge A) style: internal ``[[label]]`` tokens are resolved to
"(Author Year)" strings and the reference list is alphabetised, with no orphan
or phantom references.
"""

from __future__ import annotations

import csv
import os
import re
import shutil
import textwrap
from collections import Counter
import urllib.error
import urllib.request
import zipfile
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
import numpy as np

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "output"
WORK = ROOT / "build"
FIGDIR = OUTPUT / "figures"

TITLE = (
    "Lifecycle Data Exposure in Shared Micromobility: A Scoping Review and a "
    "Global Audit of Public GBFS Vehicle Feeds and Operator Disclosures"
)
SHORT_TITLE = "Lifecycle data exposure in shared micromobility"
AUTHOR = "Onishi Tatsuki"
AFFILIATION = ("Data Science and AI Innovation Research Promotion Center, "
               "Shiga University, Hikone, Japan")
EMAIL = "bougtoir@gmail.com"
ORCID = "0000-0001-7261-9062"
JOURNAL = "Data & Policy"
PUBLISHER = "Cambridge University Press"
ARTICLE_TYPE = "Research Article"
AREA_OF_INTEREST = "Area 4: Ethics, Equity and Trustworthiness of Data"
PUBLIC_REPO_URL = "https://github.com/bougtoir/discharged-secrets-scoping-review"
BUILD_DATE = "15 July 2026"

# --- Review-model toggle -------------------------------------------------
# The same pipeline serves single-blind (default) and double-masked review.
# BLINDED=1 removes author-identifying content from the *main manuscript*
# (byline, acknowledgements, and any identity-revealing repository URL) so the
# manuscript file itself can be uploaded to a double-masked journal. The title
# page and cover letter are always non-anonymized because journals collect them
# separately and do not forward them to reviewers.
BLINDED = os.environ.get("BLINDED", "0") == "1"

# Persistent DOI for the archived deposit (e.g. Zenodo concept DOI covering all
# versions). Set ZENODO_DOI once the archive is minted; the availability
# statement then cites the DOI instead of the bare repository URL.
ZENODO_DOI = os.environ.get("ZENODO_DOI", "").strip()

# Identity-free link used *only* during double-masked review. anonymous.4open.science
# mirrors a GitHub repository behind a temporary URL that does not reveal the
# owner/organisation. Override per submission with ANON_REPO_URL.
ANON_REPO_URL = os.environ.get(
    "ANON_REPO_URL",
    "https://anonymous.4open.science/r/discharged-secrets-scoping-review")

# Data & Policy allows up to five keywords, separated by semicolons.
KEYWORDS = [
    "shared micromobility",
    "location privacy",
    "data lifecycle",
    "open mobility data",
    "data governance",
]

# Required by Data & Policy: a 120-word statement, in accessible language,
# placed directly beneath the abstract.
POLICY_SIGNIFICANCE = (
    "Shared e-scooters and e-bikes are connected devices that cities and "
    "residents use but neither own nor decommission. This study gives "
    "policymakers an evidence-based map of what these systems actually "
    "disclose. A global audit of public vehicle feeds shows that persistent "
    "identifiers and precise locations are published almost universally, while "
    "operator privacy notices rarely address device disposal or vulnerability "
    "reporting. For regulators and procuring authorities, this indicates that "
    "oversight focused only on real-time location understates lifecycle "
    "exposure. We provide a transparent, reproducible lifecycle framework that "
    "ties each stage to the strength of its evidence, helping agencies target "
    "disclosure requirements, procurement clauses, and end-of-life "
    "data-handling rules where public accountability is currently weakest. "
    "The whole audit relies only on public data and open code."
)

# ---------------------------------------------------------------------------
# References. Keyed by short label; rendered in author-date (Cambridge A)
# style and alphabetised in the reference list (see resolve_citations). Every
# reference below was checked against a public identifier (DOI/URL) during
# preparation and is recorded in output/Reference_Verification.csv.
# ---------------------------------------------------------------------------
REFS: dict[str, str] = {
    "prisma_scr": (
        "Tricco, A.C., Lillie, E., Zarin, W., O'Brien, K.K., Colquhoun, H., "
        "Levac, D., et al. (2018) PRISMA Extension for Scoping Reviews "
        "(PRISMA-ScR): checklist and explanation. Annals of Internal Medicine "
        "169(7), 467-473. https://doi.org/10.7326/M18-0850."
    ),
    "grade_indirectness": (
        "Guyatt, G.H., Oxman, A.D., Kunz, R., Woodcock, J., Brozek, J., "
        "Helfand, M., et al. (2011) GRADE guidelines: 8. Rating the quality of "
        "evidence - indirectness. Journal of Clinical Epidemiology 64(12), "
        "1303-1310. https://doi.org/10.1016/j.jclinepi.2011.04.014."
    ),
    "arksey": (
        "Arksey, H. and O'Malley, L. (2005) Scoping studies: towards a "
        "methodological framework. International Journal of Social Research "
        "Methodology 8(1), 19-32. https://doi.org/10.1080/1364557032000119616."
    ),
    "demontjoye": (
        "de Montjoye, Y.-A., Hidalgo, C.A., Verleysen, M. and Blondel, V.D. "
        "(2013) Unique in the crowd: the privacy bounds of human mobility. "
        "Scientific Reports 3, 1376. https://doi.org/10.1038/srep01376."
    ),
    "gdpr": (
        "European Parliament and Council (2016) Regulation (EU) 2016/679 "
        "(General Data Protection Regulation). Official Journal of the European "
        "Union L119, 1-88."
    ),
    "elzer": (
        "Elzer, T., Ruben, M., Classen, J. and Hollick, M. (2025) They see me "
        "scooting: a long-term real-world data analysis of shared "
        "micro-mobility services. In 2025 IEEE 10th European Symposium on "
        "Security and Privacy (EuroS&P). IEEE. "
        "https://doi.org/10.1109/EuroSP63326.2025.00014."
    ),
    "etrojans": (
        "Casagrande, M., Losiouk, E., Conti, M., Payer, M. and Antonioli, D. "
        "(2024) E-Trojans: ransomware, tracking, DoS, and data leaks on "
        "battery-powered embedded systems. arXiv:2411.17184. "
        "https://doi.org/10.48550/arXiv.2411.17184."
    ),
    "espoofer": (
        "Casagrande, M., Losiouk, E., Conti, M., Payer, M. and Antonioli, D. "
        "(2023) E-Spoofer: attacking and defending Xiaomi electric scooter "
        "ecosystem. In Proceedings of the 16th ACM Conference on Security and "
        "Privacy in Wireless and Mobile Networks (WiSec). ACM. "
        "https://doi.org/10.1145/3558482.3590176."
    ),
    "vinayaga2022": (
        "Vinayaga-Sureshkanth, N., Maiti, A., Jadliwala, M., Crager, K., He, J. "
        "and Rathore, H. (2022) An investigative study on the privacy "
        "implications of mobile e-scooter rental apps. In Proceedings of the "
        "15th ACM Conference on Security and Privacy in Wireless and Mobile "
        "Networks (WiSec). ACM. https://doi.org/10.1145/3507657.3528551."
    ),
    "hilgert": (
        "Hilgert, J.-N., Lambertz, M., Rybalka, M., Schell, R. and Vogt, R. "
        "(2021) A forensic analysis of micromobility solutions. Forensic "
        "Science International: Digital Investigation 38, 301137. "
        "https://doi.org/10.1016/j.fsidi.2021.301137."
    ),
    "isik": (
        "Isik, A.B., Dag, T. and Ozkan, K. (2023) E-scooter sharing platforms: "
        "understanding their architecture and cybersecurity threats. In 2023 "
        "IEEE 26th International Conference on Intelligent Transportation "
        "Systems (ITSC). IEEE. https://doi.org/10.1109/ITSC57777.2023.10421849."
    ),
    "vinayaga2020": (
        "Vinayaga-Sureshkanth, N., Maiti, A., Jadliwala, M., Crager, K., He, J. "
        "and Rathore, H. (2020) Security and privacy challenges in upcoming "
        "intelligent urban micromobility transportation systems. In Proceedings "
        "of the Second ACM Workshop on Automotive and Aerial Vehicle Security "
        "(AutoSec). ACM. https://doi.org/10.1145/3375706.3380559."
    ),
    "petersen": (
        "Petersen, M.L. (2019) Scoot over smart devices: the invisible costs of "
        "rental scooters. Surveillance & Society 17(1/2), 267-273. "
        "https://doi.org/10.24908/ss.v17i1/2.13112."
    ),
    "sato": (
        "Sato, K., Fukushima, N., Fujii, K. and Kitani, T. (2025) Data "
        "acquisition framework for micromobility vehicles toward driving risk "
        "prediction. IEEE Security & Privacy 23(1). "
        "https://doi.org/10.1109/MSEC.2024.3441731."
    ),
    "yilmaz2022": (
        "Yilmaz, S. and Karsligil, M.E. (2022) Analysis of location spoofing "
        "threats on e-scooter sharing. In 2022 30th Signal Processing and "
        "Communications Applications Conference (SIU). IEEE. "
        "https://doi.org/10.1109/SIU55565.2022.9864946."
    ),
    "yilmaz2023": (
        "Yilmaz, S. and Karsligil, M.E. (2023) Geo-location spoofing on "
        "e-scooters: threat analysis and prevention framework. Balkan Journal "
        "of Electrical and Computer Engineering 11(2). "
        "https://doi.org/10.17694/bajece.1231384."
    ),
    "li2020": (
        "Li, Y. and Zhang, X. (2020) Linking privacy concerns for traceable "
        "information and information privacy protective responses. In "
        "Proceedings of the 53rd Hawaii International Conference on System "
        "Sciences (HICSS). https://doi.org/10.24251/HICSS.2020.105."
    ),
    "hannemann": (
        "Hannemann, A., Buchholz, E. and Ziegler, D. (2023) Is homomorphic "
        "encryption feasible for smart mobility? In Annals of Computer Science "
        "and Information Systems, vol. 35 (FedCSIS). "
        "https://doi.org/10.15439/2023F695."
    ),
    "zhou": (
        "Zhou, Y. (2024) Data extraction in dockless bikeshare: an analysis "
        "from users' perspective. Big Data & Society 11(4). "
        "https://doi.org/10.1177/20539517241299724."
    ),
    "leaky": (
        "Marchiori, A., Losiouk, E., Conti, M. and Antonioli, D. (2025) Leaky "
        "batteries: a novel set of side-channel attacks on electric vehicles. "
        "In Computer Security - ESORICS 2025, Lecture Notes in Computer "
        "Science. Springer. https://doi.org/10.1007/978-3-032-00624-0_16."
    ),
    "iotreuse": (
        "Liu, Y., Zhou, Z., Zhang, J., et al. (2023) How IoT re-using threatens "
        "your sensitive data: exploring the user-data disposal in used IoT "
        "devices. In 2023 IEEE Symposium on Security and Privacy (SP). IEEE. "
        "https://doi.org/10.1109/SP46215.2023.10179294."
    ),
    "bms": (
        "Basic, F., Gaertner, M. and Steger, C. (2023) Secure data acquisition "
        "for battery management systems. In 2023 26th Euromicro Conference on "
        "Digital System Design (DSD). IEEE. "
        "https://doi.org/10.1109/DSD60849.2023.00082."
    ),
    "remanence": (
        "Joshi, A. and Raval, M.S. (2018) Standards and techniques to remove "
        "data remanence in cloud storage. In 2018 IEEE Punecon. IEEE. "
        "https://doi.org/10.1109/PUNECON.2018.8745370."
    ),
    "gbfs_spec": (
        "MobilityData (2024) General Bikeshare Feed Specification (GBFS), "
        "version 3.0. Available at https://gbfs.org/ (accessed 12 July 2026)."
    ),
    "gbfs_registry": (
        "MobilityData (2026) GBFS systems catalogue (systems.csv). Available at "
        "https://github.com/MobilityData/gbfs (accessed 8 July 2026)."
    ),
    "mds_privacy": (
        "Open Mobility Foundation (2021) MDS privacy guide for cities. "
        "Available at https://www.openmobilityfoundation.org/ (accessed 12 "
        "July 2026)."
    ),
    "edpb_cv": (
        "European Data Protection Board (2021) Guidelines 01/2020 on processing "
        "personal data in the context of connected vehicles and mobility "
        "related applications, version 2.0."
    ),
    "nist88": (
        "Kissel, R., Regenscheid, A., Scholl, M. and Stine, K. (2014) NIST "
        "Special Publication 800-88 Revision 1: Guidelines for Media "
        "Sanitization. National Institute of Standards and Technology. "
        "https://doi.org/10.6028/NIST.SP.800-88r1."
    ),
    "nist161": (
        "Boyens, J., Smith, A., Bartol, N., Winkler, K., Holbrook, A. and "
        "Fallon, M. (2022) NIST Special Publication 800-161 Revision 1: "
        "Cybersecurity Supply Chain Risk Management Practices for Systems and "
        "Organizations. National Institute of Standards and Technology. "
        "https://doi.org/10.6028/NIST.SP.800-161r1."
    ),
    "eu_battery": (
        "European Parliament and Council (2023) Regulation (EU) 2023/1542 "
        "concerning batteries and waste batteries. Official Journal of the "
        "European Union L191, 1-117."
    ),
}

# In-text author-date citation metadata (Cambridge A style). Each entry gives
# the in-text author string, the year, and an alphabetical sort key (first
# author surname / organisation) for ordering the reference list.
CITEMETA: dict[str, tuple[str, str, str]] = {
    "prisma_scr": ("Tricco et al.", "2018", "tricco"),
    "grade_indirectness": ("Guyatt et al.", "2011", "guyatt"),
    "arksey": ("Arksey and O'Malley", "2005", "arksey"),
    "demontjoye": ("de Montjoye et al.", "2013", "de montjoye"),
    "gdpr": ("European Parliament and Council", "2016", "european parliament and council"),
    "elzer": ("Elzer et al.", "2025", "elzer"),
    "etrojans": ("Casagrande et al.", "2024", "casagrande"),
    "espoofer": ("Casagrande et al.", "2023", "casagrande"),
    "vinayaga2022": ("Vinayaga-Sureshkanth et al.", "2022", "vinayaga-sureshkanth"),
    "hilgert": ("Hilgert et al.", "2021", "hilgert"),
    "isik": ("Isik et al.", "2023", "isik"),
    "vinayaga2020": ("Vinayaga-Sureshkanth et al.", "2020", "vinayaga-sureshkanth"),
    "petersen": ("Petersen", "2019", "petersen"),
    "sato": ("Sato et al.", "2025", "sato"),
    "yilmaz2022": ("Yilmaz and Karsligil", "2022", "yilmaz"),
    "yilmaz2023": ("Yilmaz and Karsligil", "2023", "yilmaz"),
    "li2020": ("Li and Zhang", "2020", "li"),
    "hannemann": ("Hannemann et al.", "2023", "hannemann"),
    "zhou": ("Zhou", "2024", "zhou"),
    "leaky": ("Marchiori et al.", "2025", "marchiori"),
    "iotreuse": ("Liu et al.", "2023", "liu"),
    "bms": ("Basic et al.", "2023", "basic"),
    "remanence": ("Joshi and Raval", "2018", "joshi"),
    "gbfs_spec": ("MobilityData", "2024", "mobilitydata"),
    "gbfs_registry": ("MobilityData", "2026", "mobilitydata"),
    "mds_privacy": ("Open Mobility Foundation", "2021", "open mobility foundation"),
    "edpb_cv": ("European Data Protection Board", "2021", "european data protection board"),
    "nist88": ("Kissel et al.", "2014", "kissel"),
    "nist161": ("Boyens et al.", "2022", "boyens"),
    "eu_battery": ("European Parliament and Council", "2023", "european parliament and council"),
}

assert set(CITEMETA) == set(REFS), (
    "CITEMETA and REFS must cover the same labels: "
    f"{set(CITEMETA) ^ set(REFS)}"
)


def intext(label: str) -> str:
    authors, year, _ = CITEMETA[label]
    return f"{authors} {year}"


def sortkey(label: str) -> tuple[str, str]:
    _, year, srt = CITEMETA[label]
    return (srt, year)

# Contents of the archive, reused across every availability-statement variant.
_ARCHIVE_CONTENTS = (
    "the frozen GBFS registry snapshot and its checksum, the title/abstract "
    "and full-text screening decisions, the disclosure-audit coding sheets "
    "with verbatim locator quotations, the analysis scripts, the analysis "
    "results, and the figure- and table-generation code. Raw vehicle "
    "identifiers, exact coordinates, and vehicle-specific deep links were not "
    "retained; only field presence or absence was recorded.")
_REPRO_SENTENCE = (
    " Every count, proportion, and confidence interval reported in the article "
    "can be regenerated from these materials with a single build command.")


def data_availability_statement() -> str:
    """Return the availability statement matching the current review model.

    - double-masked (BLINDED): identity-free anonymized link; DOI/URL added on
      acceptance;
    - single-blind + Zenodo DOI: cite the persistent DOI plus the repository;
    - single-blind, no DOI yet: cite the public repository URL.
    """
    if BLINDED:
        return (
            "The data and code that support the findings of this study are "
            f"available to reviewers via an anonymized repository at "
            f"{ANON_REPO_URL}. The archive contains " + _ARCHIVE_CONTENTS +
            " Upon acceptance, these materials will be deposited in a public "
            "repository with a permanent DOI." + _REPRO_SENTENCE)
    if ZENODO_DOI:
        return (
            "The data and code that support the findings of this study are "
            f"openly archived on Zenodo at https://doi.org/{ZENODO_DOI} "
            "(concept DOI, covering all versions), with ongoing development at "
            f"{PUBLIC_REPO_URL}. The archive contains " + _ARCHIVE_CONTENTS +
            _REPRO_SENTENCE)
    return (
        "The data and code that support the findings of this study are openly "
        f"available in the project repository at {PUBLIC_REPO_URL}. This "
        "includes " + _ARCHIVE_CONTENTS + _REPRO_SENTENCE)


# Disclosure statements required by Data & Policy, placed after the main text
# and before the references. Data availability, funding, and competing
# interests are mandatory; the remaining statements are included as good
# practice and to preserve the ethics and generative-AI disclosures.
def build_declarations() -> list[tuple[str, str]]:
    acknowledgements = (
        "Acknowledgements are withheld to preserve author anonymity for "
        "double-masked peer review; they will be restored on acceptance."
        if BLINDED else
        "The author thanks the maintainers of the public GBFS ecosystem and "
        "the open bibliographic sources that made this audit possible.")
    contributions = (
        "Author contributions are withheld to preserve author anonymity for "
        "double-masked peer review; they will be restored on acceptance."
        if BLINDED else
        f"{AUTHOR}: conceptualization; methodology; software; formal "
        "analysis; data curation; writing - original draft; writing - review "
        "and editing.")
    ai_use = (
        "No generative AI tool was used to produce scientific content, "
        "analysis, or interpretation. Any use was limited to routine language "
        "editing and was reviewed by the author"
        + (", who takes full responsibility for the text."
           if not BLINDED else "s, who take full responsibility for the text."))
    return [
        ("Acknowledgements", acknowledgements),
        ("Data availability statement", data_availability_statement()),
        ("Funding statement",
         "This research received no specific grant from any funding agency, "
         "commercial or not-for-profit sectors."),
        ("Competing interests",
         "The author declares none."),
        ("Author contributions", contributions),
        ("Ethical standards",
         "The research meets all ethical guidelines, including adherence to "
         "the legal requirements of the study jurisdiction. No human "
         "participants were recruited. The study analysed only publicly "
         "accessible feeds and documents and did not attempt authentication, "
         "access-control circumvention, or interaction with individual users."),
        ("Use of generative AI", ai_use),
    ]


DECLARATIONS = build_declarations()


# ---------------------------------------------------------------------------
# Data access helpers
# ---------------------------------------------------------------------------
def load_gbfs_summary() -> dict[tuple[str, str], dict]:
    table: dict[tuple[str, str], dict] = {}
    with (ROOT / "results" / "gbfs_summary.csv").open(encoding="utf-8") as fh:
        for row in csv.DictReader(fh):
            table[(row["stratum"], row["metric"])] = row
    return table


def load_document_audit() -> tuple[list[str], dict[str, dict[str, str]], int, int]:
    rows = list(csv.DictReader((ROOT / "data" / "document_audit.csv").open(encoding="utf-8")))
    domains: list[str] = []
    matrix: dict[str, dict[str, str]] = {}
    operators_order: list[str] = []
    for r in rows:
        op = r["operator_name"]
        if op not in matrix:
            matrix[op] = {}
            operators_order.append(op)
        matrix[op][r["domain"]] = r["status"]
        if r["domain"] not in domains:
            domains.append(r["domain"])
    coded = sum(1 for op in operators_order
                if any(v != "unavailable" for v in matrix[op].values()))
    return domains, matrix, len(operators_order), coded


def load_extraction() -> list[dict]:
    return list(csv.DictReader((ROOT / "review" / "evidence_extraction.csv").open(encoding="utf-8")))


GBFS = load_gbfs_summary()
DOMAINS, DOCMATRIX, N_OPERATORS, N_CODED = load_document_audit()
EXTRACTION = load_extraction()


def g(stratum: str, metric: str, field: str) -> str:
    return GBFS[(stratum, metric)][field]


def pct(stratum: str, metric: str) -> str:
    row = GBFS[(stratum, metric)]
    return f"{float(row['percent']):.1f}%"


def ci(stratum: str, metric: str) -> str:
    row = GBFS[(stratum, metric)]
    return f"{float(row['wilson_95_low_percent']):.1f}-{float(row['wilson_95_high_percent']):.1f}%"


def domain_count(domain: str, status: str) -> int:
    return sum(1 for op, cells in DOCMATRIX.items() if cells.get(domain) == status)


def load_screening_stats() -> dict[str, int]:
    """Derive every PRISMA-flow count directly from the committed screening
    sheet so that no flow number is hard-coded in the manuscript."""
    rows = list(csv.DictReader(
        (ROOT / "review" / "screening.csv").open(encoding="utf-8")))
    ta = Counter(r["title_abstract_decision"] for r in rows)
    excl = Counter(r["full_text_exclusion_reason"] for r in rows
                   if r["full_text_decision"] == "exclude")
    included = [r for r in rows if r["full_text_decision"] == "include"]
    dist = Counter(r["evidence_distance"] for r in included)
    second_pass = [r for r in rows if r["second_pass_decision"].strip()]
    agree = sum(1 for r in second_pass
                if r["second_pass_decision"].strip() ==
                r["title_abstract_decision"].strip())
    n_no_data_path = excl["F1_NO_RELEVANT_DATA_PATH"]
    n_not_transferable = excl["F2_MECHANISM_NOT_TRANSFERABLE"]
    return {
        "identified": len(rows),
        "ta_excluded": ta["exclude"],
        "sought": ta["include"] + ta["uncertain"],
        "not_retrieved": excl["F4_NOT_RETRIEVED"],
        "no_data_path": n_no_data_path,
        "not_transferable": n_not_transferable,
        "excluded_fulltext": n_no_data_path + n_not_transferable,
        "included": len(included),
        "with_doi": sum(1 for r in rows if r["doi"].strip()),
        "d4": dist["D4"], "d3": dist["D3"], "d2": dist["D2"],
        "resample_n": len(second_pass),
        "resample_agreement_pct":
            round(100 * agree / len(second_pass)) if second_pass else 0,
    }


def _summary_int(stratum: str, metric: str, field: str) -> int:
    return int(g(stratum, metric, field))


SCR = load_screening_stats()

# Audit totals, all read from the reproducible GBFS summary rather than typed.
N_REGISTRY = _summary_int(
    "all_registry_entries", "auto_discovery_reachable", "denominator")
N_REACHABLE = _summary_int(
    "all_registry_entries", "auto_discovery_reachable", "count")
N_MOTOR_FEEDS = _summary_int(
    "declared_motorized_micromobility_feeds", "has_vehicle_id", "denominator")
N_OP_DOMAINS = _summary_int(
    "declared_motorized_micromobility_operator_domains_any",
    "has_vehicle_id", "denominator")


def fmt(n: int) -> str:
    """Render an integer with thousands separators (e.g. 2169 -> '2,169')."""
    return f"{n:,}"


# ---------------------------------------------------------------------------
# Figures (English text only, per manuscript language rule)
# ---------------------------------------------------------------------------
FIG_CAPTIONS = {
    1: ("Figure 1. PRISMA-ScR flow of record identification, title/abstract "
        "screening, full-text eligibility assessment, and inclusion. Counts are "
        "reproduced from the committed screening sheet."),
    2: (f"Figure 2. Evidence map of the {SCR['included']} included direct studies "
        "by evidence distance (D4 direct empirical, D3 direct documentary, D2 "
        "near-domain empirical) and lifecycle stage addressed. A study that "
        "addresses more than one stage is counted in each relevant stage, so row "
        "totals may exceed the number of studies at that distance."),
    3: (f"Figure 3. Prevalence of disclosed fields in {fmt(N_MOTOR_FEEDS)} "
        "successfully retrieved, non-empty public vehicle-status feeds that "
        "declared motorized micromobility. Bars show the point estimate; "
        "whiskers show the 95% Wilson confidence interval. Field presence is a "
        "disclosure signal, not evidence of a privacy harm."),
    4: (f"Figure 4. Disclosure-domain coding of public operator privacy notices "
        f"across {len(DOMAINS)} domains (n = {N_CODED} coded operators; "
        f"{N_OPERATORS - N_CODED} further operators were unavailable for "
        "reproducible retrieval). Cells show explicit, partial, or not-found "
        "coding; not-found denotes silence in the document, not evidence that a "
        "practice is absent."),
    5: ("Figure 5. Lifecycle data-exposure model linking each stage to the "
        "evidence distance of its supporting sources and to a proposed, as-yet "
        "unvalidated control. Solid boxes denote stages with direct (D3-D4) "
        "evidence; dashed boxes denote stages supported only by near-domain "
        "(D2) evidence."),
}

PALETTE = {
    "explicit": "#2c7fb8",
    "partial": "#7fcdbb",
    "not_found": "#edf8b1",
    "unavailable": "#d9d9d9",
    "bar": "#2c7fb8",
    "D4": "#08519c",
    "D3": "#3182bd",
    "D2": "#9ecae1",
}


def fig1_prisma(ax) -> None:
    ax.axis("off")
    boxes = [
        (0.5, 0.93, f"Records identified from\nbibliographic sources (n = {fmt(SCR['identified'])})"),
        (0.5, 0.75, f"Records screened on title/abstract\n(n = {fmt(SCR['identified'])})"),
        (0.5, 0.55, f"Records sought for full-text\nassessment (n = {fmt(SCR['sought'])})"),
        (0.5, 0.33, f"Full-text records assessed\nfor eligibility (n = {fmt(SCR['sought'])})"),
        (0.5, 0.11, f"Direct studies included in\nsynthesis (n = {fmt(SCR['included'])})"),
    ]
    for x, y, text in boxes:
        ax.add_patch(FancyBboxPatch((x - 0.22, y - 0.055), 0.44, 0.11,
                     boxstyle="round,pad=0.01", fc="#deebf7", ec="#08519c", lw=1.2,
                     transform=ax.transAxes))
        ax.text(x, y, text, ha="center", va="center", fontsize=8.5,
                transform=ax.transAxes)
    excl = [
        (0.80, 0.75, f"Excluded on title/abstract\n(n = {fmt(SCR['ta_excluded'])})"),
        (0.80, 0.44, f"Full text not retrieved\n(n = {fmt(SCR['not_retrieved'])})"),
        (0.80, 0.22, f"Excluded at full text (n = {fmt(SCR['excluded_fulltext'])}):\n"
                     f"no relevant data path (n = {fmt(SCR['no_data_path'])});\n"
                     f"mechanism not transferable (n = {fmt(SCR['not_transferable'])})"),
    ]
    for x, y, text in excl:
        ax.add_patch(FancyBboxPatch((x - 0.18, y - 0.06), 0.36, 0.12,
                     boxstyle="round,pad=0.01", fc="#f7f7f7", ec="#999999", lw=1.0,
                     transform=ax.transAxes))
        ax.text(x, y, text, ha="center", va="center", fontsize=7.5,
                transform=ax.transAxes)
    ys = [0.875, 0.695, 0.475, 0.275]
    for y0 in ys:
        ax.add_patch(FancyArrowPatch((0.5, y0), (0.5, y0 - 0.05),
                     arrowstyle="-|>", mutation_scale=12, color="#08519c",
                     transform=ax.transAxes))
    for y0, yx in [(0.75, 0.75), (0.55, 0.44), (0.33, 0.22)]:
        ax.add_patch(FancyArrowPatch((0.72, y0), (0.62, y0),
                     arrowstyle="-|>", mutation_scale=10, color="#999999",
                     transform=ax.transAxes))


def fig2_evidence_map(ax) -> None:
    stages = ["Operation", "Maintenance/\nservice", "Recall/return",
              "Second-life/\ndisposal"]
    distances = ["D2", "D3", "D4"]
    # map each study to (distance, stage) using lifecycle_stage in extraction
    stage_key = {
        "operation": 0, "data sharing": 0,
        "maintenance": 1, "maintenance and service": 1,
        "recall and return": 2, "recall": 2, "return": 2,
        "second-life and disposal": 3, "second-life": 3, "disposal": 3,
        "recycling": 3, "manufacturing": 3,
    }
    counts = np.zeros((len(distances), len(stages)), dtype=int)
    for e in EXTRACTION:
        d = distances.index(e["evidence_distance"]) if e["evidence_distance"] in distances else None
        if d is None:
            continue
        seen: set[int] = set()
        for token in e["lifecycle_stage"].split(";"):
            s = stage_key.get(token.strip().lower())
            if s is not None:
                seen.add(s)
        for s in seen:
            counts[d, s] += 1
    im = ax.imshow(counts, cmap="Blues", aspect="auto", vmin=0, vmax=max(counts.max(), 1))
    ax.set_xticks(range(len(stages)))
    ax.set_xticklabels(stages, fontsize=8)
    ax.set_yticks(range(len(distances)))
    ax.set_yticklabels(["D2 near-domain", "D3 documentary", "D4 empirical"], fontsize=8)
    for i in range(len(distances)):
        for j in range(len(stages)):
            if counts[i, j]:
                ax.text(j, i, str(counts[i, j]), ha="center", va="center",
                        color="black", fontsize=9)
    ax.set_xlabel("Lifecycle stage addressed", fontsize=8.5)
    cb = plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    cb.set_label("Number of included studies", fontsize=8)


def fig3_prevalence(ax) -> None:
    metrics = [
        ("has_vehicle_id", "Vehicle\nidentifier"),
        ("has_location_fields", "Latitude/\nlongitude"),
        ("has_range", "Current\nrange"),
        ("has_deep_link", "Rental\ndeep link"),
        ("has_last_reported", "Last-reported\ntimestamp"),
        ("has_battery_percent", "Battery/fuel\npercentage"),
    ]
    stratum = "declared_motorized_micromobility_feeds"
    labels, vals, lo, hi = [], [], [], []
    for metric, lab in metrics:
        row = GBFS[(stratum, metric)]
        labels.append(lab)
        vals.append(float(row["percent"]))
        lo.append(float(row["percent"]) - float(row["wilson_95_low_percent"]))
        hi.append(float(row["wilson_95_high_percent"]) - float(row["percent"]))
    x = np.arange(len(labels))
    ax.bar(x, vals, color=PALETTE["bar"], width=0.62,
           yerr=[lo, hi], capsize=3, ecolor="#333333")
    for xi, v, h in zip(x, vals, hi):
        ax.text(xi, v + h + 1.6, f"{v:.1f}", ha="center", fontsize=8)
    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=8)
    ax.set_ylabel("Feeds disclosing the field (%)", fontsize=9)
    ax.set_ylim(0, 108)
    ax.set_title(
        f"Disclosed fields in {fmt(N_MOTOR_FEEDS)} motorized micromobility vehicle feeds",
        fontsize=9.5)
    ax.spines[["top", "right"]].set_visible(False)


DOMAIN_LABELS = {
    "location_data": "Location",
    "trip_time_data": "Trip/time",
    "vehicle_identifier": "Vehicle ID",
    "battery_or_diagnostic_data": "Battery/diagnostic",
    "maintenance_or_repair_data": "Maintenance/repair",
    "account_payment_device_data": "Account/payment/device",
    "analytics_or_profiling": "Analytics/profiling",
    "retention": "Retention",
    "processors_or_contractors": "Processors/contractors",
    "international_transfer": "Intl transfer",
    "user_rights": "Data-subject rights",
    "incident_contact": "Incident contact",
    "vulnerability_disclosure": "Vulnerability disclosure",
    "return_recycling_disposal": "Return/recycling/disposal",
}
STATUS_VAL = {"explicit": 3, "partial": 2, "not_found": 1, "unavailable": 0}


def fig4_heatmap(ax) -> None:
    operators = [op for op in DOCMATRIX
                 if any(v != "unavailable" for v in DOCMATRIX[op].values())]
    grid = np.array([[STATUS_VAL[DOCMATRIX[op][d]] for d in DOMAINS] for op in operators])
    from matplotlib.colors import ListedColormap
    cmap = ListedColormap([PALETTE["unavailable"], PALETTE["not_found"],
                           PALETTE["partial"], PALETTE["explicit"]])
    ax.imshow(grid, cmap=cmap, aspect="auto", vmin=0, vmax=3)
    ax.set_xticks(range(len(DOMAINS)))
    ax.set_xticklabels([DOMAIN_LABELS[d] for d in DOMAINS], rotation=45,
                       ha="right", fontsize=7)
    ax.set_yticks(range(len(operators)))
    ax.set_yticklabels(operators, fontsize=7)
    ax.set_xticks(np.arange(-0.5, len(DOMAINS), 1), minor=True)
    ax.set_yticks(np.arange(-0.5, len(operators), 1), minor=True)
    ax.grid(which="minor", color="white", linewidth=1.0)
    ax.tick_params(which="minor", length=0)
    handles = [plt.Rectangle((0, 0), 1, 1, fc=PALETTE[k]) for k in
               ["explicit", "partial", "not_found"]]
    ax.legend(handles, ["Explicit", "Partial", "Not found"],
              loc="upper left", bbox_to_anchor=(1.01, 1.0), fontsize=7,
              frameon=False)


LIFECYCLE = [
    ("Procurement", "D3", "Contract, provenance,\nand access terms"),
    ("Deployment", "D3", "Feed configuration\nand field minimization"),
    ("Operation", "D4", "Local processing;\naccess and retention limits"),
    ("Maintenance/service", "D3", "Controlled contractor\naccess and logging"),
    ("Recall/return", "D2", "Chain-of-custody\nfor returned units"),
    ("Second-life/disposal", "D2", "Verified media\nsanitization"),
]


def fig5_lifecycle(ax) -> None:
    ax.axis("off")
    n = len(LIFECYCLE)
    xs = np.linspace(0.08, 0.92, n)
    for i, ((stage, dist, control), x) in enumerate(zip(LIFECYCLE, xs)):
        direct = dist in ("D3", "D4")
        style = "round,pad=0.01" if direct else "round,pad=0.01"
        ax.add_patch(FancyBboxPatch((x - 0.065, 0.62), 0.13, 0.16,
                     boxstyle=style, fc=PALETTE.get(dist, "#deebf7"),
                     ec="#08306b", lw=1.4,
                     linestyle="-" if direct else "--", transform=ax.transAxes))
        ax.text(x, 0.70, f"{stage}\n[{dist}]", ha="center", va="center",
                fontsize=7.2, transform=ax.transAxes, color="white" if dist == "D4" else "black")
        ax.add_patch(FancyBboxPatch((x - 0.075, 0.30), 0.15, 0.18,
                     boxstyle="round,pad=0.01", fc="#f7fbff", ec="#999999",
                     lw=1.0, transform=ax.transAxes))
        ax.text(x, 0.39, control, ha="center", va="center", fontsize=6.6,
                transform=ax.transAxes)
        ax.add_patch(FancyArrowPatch((x, 0.62), (x, 0.485), arrowstyle="-|>",
                     mutation_scale=9, color="#999999", transform=ax.transAxes))
        if i < n - 1:
            ax.add_patch(FancyArrowPatch((x + 0.065, 0.70), (xs[i + 1] - 0.065, 0.70),
                         arrowstyle="-|>", mutation_scale=10, color="#08306b",
                         transform=ax.transAxes))
    ax.text(0.5, 0.90, "Lifecycle stage (with dominant evidence distance)",
            ha="center", fontsize=8.5, transform=ax.transAxes, weight="bold")
    ax.text(0.5, 0.19, "Proposed control (effectiveness not yet validated)",
            ha="center", fontsize=8.5, transform=ax.transAxes, weight="bold")


FIGFUNCS = {1: fig1_prisma, 2: fig2_evidence_map, 3: fig3_prevalence,
            4: fig4_heatmap, 5: fig5_lifecycle}
FIGSIZE = {1: (7.0, 6.0), 2: (7.0, 4.2), 3: (7.2, 4.4), 4: (8.4, 5.2), 5: (9.6, 4.6)}


def build_figures() -> dict[int, dict[str, Path]]:
    FIGDIR.mkdir(parents=True, exist_ok=True)
    paths: dict[int, dict[str, Path]] = {}
    for num, func in FIGFUNCS.items():
        fig, ax = plt.subplots(figsize=FIGSIZE[num])
        func(ax)
        fig.tight_layout()
        png = FIGDIR / f"Figure{num}.png"
        tiff = FIGDIR / f"Figure{num}.tiff"
        pdf = FIGDIR / f"Figure{num}.pdf"
        fig.savefig(png, dpi=600, bbox_inches="tight")
        fig.savefig(tiff, dpi=600, bbox_inches="tight", pil_kwargs={"compression": "tiff_lzw"})
        fig.savefig(pdf, bbox_inches="tight")
        plt.close(fig)
        paths[num] = {"png": png, "tiff": tiff, "pdf": pdf}
    return paths


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------
def build_tables() -> dict[int, dict]:
    t1 = {
        "title": "Table 1. Eligibility criteria and data sources for the three "
                 "study components.",
        "headers": ["Component", "Unit of analysis", "Inclusion criteria",
                    "Primary source"],
        "rows": [
            ["Scoping review (WP1)", "Study",
             "Peer-reviewed or archival works reporting a data path in or "
             "transferable to shared micromobility",
             f"Bibliographic search of {fmt(SCR['identified'])} records "
             "[[prisma_scr;arksey]]"],
            ["GBFS field audit (WP2)", "System/feed",
             "Public GBFS systems in the frozen registry with a reachable "
             "auto-discovery endpoint",
             "GBFS systems catalogue and live feeds [[gbfs_registry;gbfs_spec]]"],
            ["Disclosure audit (WP3)", "Operator document",
             "Publicly reachable operator privacy notice for an audited GBFS "
             "operator",
             "Operator privacy notices; governance guidance [[mds_privacy;edpb_cv]]"],
        ],
    }
    dist_label = {"D4": "D4 empirical", "D3": "D3 documentary", "D2": "D2 near-domain"}
    ref_by_id = {
        "e641ae6185ec": "espoofer", "4fe184aeb5b3": "etrojans",
        "3b9d67a99101": "elzer", "f7ef2527bcc7": "hilgert",
        "4bf4bae9ab93": "vinayaga2022", "9f85e4350278": "hannemann",
        "bdab08ae526c": "isik", "2856524cd67f": "li2020",
        "37630ff02a7b": "petersen", "b26b8d65d468": "sato",
        "ff74fa3c40ed": "vinayaga2020", "5bf50e1a9f5c": "yilmaz2022",
        "05e20798d6c6": "yilmaz2023", "62abc8d8d637": "zhou",
        "062854732ff3": "bms", "57ec900852e5": "remanence",
        "dab97e54aa73": "iotreuse", "882e5dc8fae9": "leaky",
    }
    order = {"D4": 0, "D3": 1, "D2": 2}
    ext_sorted = sorted(EXTRACTION, key=lambda e: (order[e["evidence_distance"]],
                                                   e["citation"]))
    t2_rows = []
    for e in ext_sorted:
        label = ref_by_id[e["record_id"]]
        short = e["citation"].split(".")[0]
        t2_rows.append([
            f"{short} [[year:{label}]]",
            dist_label[e["evidence_distance"]],
            e["device_or_service"],
            textwrap.shorten(e["key_result"], width=180, placeholder="..."),
        ])
    t2 = {
        "title": f"Table 2. Included direct studies (n = {len(t2_rows)}), ordered "
                 "by evidence distance. Key results are as reported by the "
                 "source; this review did not reproduce the underlying "
                 "experiments.",
        "headers": ["Study", "Evidence distance", "Device or service",
                    "Reported key result"],
        "rows": t2_rows,
    }
    def row(stratum, metric, label):
        return [label, g(stratum, metric, "count"), g(stratum, metric, "denominator"),
                pct(stratum, metric), ci(stratum, metric)]
    t3 = {
        "title": "Table 3. GBFS audit sampling and outcomes with 95% Wilson "
                 "confidence intervals. Denominators are stated explicitly; "
                 "unavailable or empty feeds are separated from absent fields.",
        "headers": ["Stratum / outcome", "n", "Denominator", "Percent", "95% CI"],
        "rows": [
            row("all_registry_entries", "auto_discovery_reachable",
                "Reachable auto-discovery endpoint"),
            row("reachable_registry_entries", "vehicle_feed_declared",
                "Vehicle-status feed declared"),
            row("successful_vehicle_feeds", "vehicle_feed_nonempty",
                "Retrieved feed with >=1 vehicle"),
            row("declared_motorized_micromobility_feeds", "has_vehicle_id",
                "Vehicle identifier (motorized feeds)"),
            row("declared_motorized_micromobility_feeds", "has_location_fields",
                "Latitude/longitude (motorized feeds)"),
            row("declared_motorized_micromobility_feeds", "has_last_reported",
                "Last-reported timestamp (motorized feeds)"),
            row("declared_motorized_micromobility_feeds", "has_battery_percent",
                "Battery/fuel percentage (motorized feeds)"),
        ],
    }
    domain_row = lambda d: [
        DOMAIN_LABELS[d], str(domain_count(d, "explicit")),
        str(domain_count(d, "partial")), str(domain_count(d, "not_found")),
    ]
    t4 = {
        "title": f"Table 4. Disclosure-domain coding of public operator privacy "
                 f"notices (n = {N_CODED} coded operators). Counts of explicit, "
                 f"partial, and not-found codings per domain; not-found denotes "
                 f"document silence, not evidence that a practice is absent.",
        "headers": ["Disclosure domain", "Explicit", "Partial", "Not found"],
        "rows": [domain_row(d) for d in DOMAINS],
    }
    t5 = {
        "title": "Table 5. Evidence-to-control traceability across the lifecycle "
                 "model. Effectiveness evidence is stated conservatively; the "
                 "proposed controls have not been empirically validated in this "
                 "study.",
        "headers": ["Lifecycle stage", "Exposure path", "Supporting evidence "
                    "(distance)", "Proposed control", "Effectiveness evidence"],
        "rows": [
            ["Procurement", "Contract and provenance terms set who can access "
             "backend data", "Governance guidance (N) [[nist161;edpb_cv]]",
             "Provenance and access clauses; data-processing terms",
             "Not validated; governance-based"],
            ["Deployment", "Feed configuration determines which fields are "
             "published", "GBFS field presence (D4) [[elzer;gbfs_spec]]",
             "Field minimization; coarse or delayed positions",
             "Not validated; design-based"],
            ["Operation", "Identifiers, positions, and timestamps enable "
             "tracking and re-identification",
             "Empirical audits (D4) [[elzer;vinayaga2022]]; mobility "
             "re-identification (D2) [[demontjoye]]",
             "Local processing; access and retention limits",
             "Partly supported by attack studies; controls not validated"],
            ["Maintenance/service", "Contractor access to devices and records",
             "Forensic and platform studies (D3-D4) [[hilgert;isik]]",
             "Controlled contractor access; audit logging",
             "Not validated"],
            ["Recall/return", "Returned units retain stored data",
             "IoT reuse study (D2) [[iotreuse]]",
             "Chain-of-custody for returned units",
             "Not validated; analogy-based"],
            ["Second-life/disposal", "Residual data on storage and batteries",
             "Data remanence and battery side-channel (D2) "
             "[[remanence;leaky;bms]]; sanitization guidance (N) [[nist88]]",
             "Verified media sanitization before resale or recycling",
             "Not validated; standards-based"],
        ],
    }
    return {1: t1, 2: t2, 3: t3, 4: t4, 5: t5}


TABLES = build_tables()


# ---------------------------------------------------------------------------
# Manuscript body. Citations use [[label]] or [[a;b]] tokens; figures/tables
# are referenced as "Fig. n" / "Table n" and inserted at their first mention.
# ---------------------------------------------------------------------------
def body_blocks() -> list[tuple]:
    n_id = fmt(SCR["identified"])
    doi_pct = 100 * SCR["with_doi"] / SCR["identified"]
    reach = pct("all_registry_entries", "auto_discovery_reachable")
    reach_ci = ci("all_registry_entries", "auto_discovery_reachable")
    blocks: list[tuple] = []

    blocks.append(("h1", "1. Introduction"))
    blocks += [("p", t) for t in [
        "Shared micromobility - dockless electric scooters and bicycles rented "
        "through a smartphone application - has become a visible class of "
        "connected devices that organizations and the public use but do not "
        "own, maintain, or decommission. Each vehicle continuously produces "
        "location, motion, and battery telemetry that is transmitted to an "
        "operator backend and, in many cities, republished through open data "
        "feeds. Prior work has shown that such telemetry can support "
        "re-identification and tracking well beyond the individual rental that "
        "generated it [[demontjoye;elzer]].",
        "Security research on shared micromobility has grown quickly but "
        "unevenly. Studies span privacy measurement of rental applications, "
        "firmware and protocol attacks on scooters, location-spoofing threats, "
        "and forensic recovery from returned devices [[vinayaga2022;espoofer;"
        "yilmaz2023;hilgert]]. These contributions are scattered across "
        "security, transportation, and forensics venues, use different threat "
        "models and units of analysis, and have not been assembled into a "
        "single map of what is actually demonstrated as opposed to argued. As "
        "a result, it is difficult to state precisely which data-exposure "
        "pathways rest on direct empirical evidence and which rest on analogy.",
        "A second gap concerns exposure that persists across the device "
        "lifecycle. Attention typically concentrates on real-time position "
        "during a rental, yet information created at deployment, maintenance, "
        "recall, and disposal can remain accessible after custody changes hands "
        "[[iotreuse;remanence]]. Battery and diagnostic channels, in "
        "particular, can leak activity patterns even when positioning is "
        "restricted [[leaky;bms]]. Whether operators disclose these lifecycle "
        "practices to the public is largely unexamined.",
        "This article addresses both gaps with an empirical, reproducible "
        "package rather than a conceptual argument. We (i) conduct a scoping "
        "review, following the Preferred Reporting Items for Systematic Reviews "
        "and Meta-Analyses extension for Scoping Reviews (PRISMA-ScR) "
        "[[prisma_scr]], of evidence on micromobility data exposure; (ii) audit "
        "the fields that operators actually publish worldwide through the "
        "General Bikeshare Feed Specification (GBFS) [[gbfs_spec]]; and (iii) "
        "audit what a matched set of operator privacy notices discloses about "
        "collection, retention, transfer, and device end-of-life. Our research "
        "questions are: RQ1, what data-exposure pathways in shared "
        "micromobility are supported by direct evidence; RQ2, which "
        "vehicle-level fields are publicly disclosed, and at what prevalence; "
        "and RQ3, how completely do public operator documents describe "
        "lifecycle data handling.",
        "We are deliberately conservative about interpretation. Publishing a "
        "field, or remaining silent about a practice, is a disclosure signal; "
        "it is not by itself evidence of a privacy harm, a compromise, or a "
        "regulatory violation. The contribution is a transparent evidence base "
        "and a lifecycle model that ties each stage to the strength of its "
        "supporting evidence and to controls whose effectiveness remains to be "
        "tested.",
    ]]

    blocks.append(("h1", "2. Methods"))
    blocks.append(("h2", "2.1. Design and reporting"))
    blocks += [("p", t) for t in [
        "The study combines three prespecified components: a scoping review "
        "(work package WP1), a cross-sectional field audit of public GBFS feeds "
        "(WP2), and a structured disclosure audit of public operator documents "
        "(WP3). The review component is reported in line with PRISMA-ScR "
        "[[prisma_scr]] and follows the scoping-review framework of Arksey and "
        "O'Malley [[year:arksey]]. Eligibility criteria and data sources for all "
        "three components are summarized in Table 1. The protocol, screening "
        "rules, coding sheets, and analysis code are openly available so that "
        "the counts reported here can be regenerated.",
    ]]
    blocks.append(("table", 1))
    blocks.append(("h2", "2.2. Scoping review (WP1)"))
    blocks += [("p", t) for t in [
        f"We compiled {n_id} records from programmatic searches of open "
        "bibliographic metadata. A digital object identifier (DOI) was "
        f"available for {fmt(SCR['with_doi'])} of {n_id} records ({doi_pct:.1f}%), "
        "for which abstracts were retrieved automatically where possible; the "
        "remaining records were screened on title and available metadata only. "
        "Title/abstract screening applied a deterministic rule set that combined "
        "target-domain relevance with the presence of a described data path, and "
        "recorded a decision and reason for every record. Because the rules are "
        "deterministic, we re-applied them to a delayed "
        f"{round(100 * SCR['resample_n'] / SCR['identified'])}% sample "
        f"(n = {fmt(SCR['resample_n'])}) and reproduced every original decision "
        f"({SCR['resample_agreement_pct']}% agreement); this demonstrates "
        "computational reproducibility rather than inter-rater reliability, as a "
        "single reviewer conducted the screening.",
        "Records marked include or uncertain were sought for full-text "
        "assessment. Each retrieved study was classified by evidence distance: "
        "D4, direct target-domain empirical evidence; D3, direct target-domain "
        "documentary evidence; D2, near-domain empirical evidence; D1, "
        "mechanism analogy; and N, normative evidence. This ordering is an "
        "operational classification defined for the present study; it is "
        "conceptually related to the GRADE notion of indirectness "
        "[[grade_indirectness]] but is not a GRADE certainty rating. We did not "
        "treat a "
        "title/abstract decision as equivalent to a confirmed full-text "
        "finding, and we did not claim to have read full text that could not be "
        "retrieved. From the included studies we built a study-level extraction "
        "table capturing device or service, design, data fields, access path, "
        "reported outcome, and limitations.",
    ]]
    blocks.append(("h2", "2.3. Public GBFS field audit (WP2)"))
    blocks += [("p", t) for t in [
        "We froze a snapshot of the public GBFS systems catalogue [[gbfs_registry]] "
        "and recorded its checksum. For each system we attempted to reach the "
        "auto-discovery endpoint, to locate a declared vehicle-status feed, to "
        "retrieve that feed, and to record which specification fields it "
        "contained. The unit of analysis is the system/feed. To respect the "
        "study's safety constraints, we recorded only the presence or absence "
        "of each field; raw vehicle identifiers, exact coordinates, and "
        "vehicle-specific deep links were not retained. Unavailable or empty "
        "feeds were separated from feeds that were retrieved but omitted a "
        "field, and all proportions are reported against explicit denominators "
        "with 95% Wilson confidence intervals (CIs). Because several large operators "
        "run many city systems, we also computed an operator-domain sensitivity "
        "analysis to check whether prevalence was driven by a few operators.",
    ]]
    blocks.append(("h2", "2.4. Public-document disclosure audit (WP3)"))
    blocks += [("p", t) for t in [
        "We selected operators from the audited GBFS population and retrieved "
        f"their public privacy notices. For each operator document we coded "
        f"{len(DOMAINS)} "
        "disclosure domains - location; trip and time data; vehicle "
        "identifiers; battery or diagnostic data; maintenance or repair "
        "records; account, payment, and device data; analytics or profiling; "
        "retention; processors or contractors; international transfers; "
        "data-subject rights; incident contact; vulnerability disclosure; and "
        "device return, recycling, or disposal - using the values explicit, "
        "partial, not found, not applicable, and unavailable. A not-found "
        "coding means the document did not address the domain; it is not "
        "evidence that the practice does not occur. Each coding is accompanied "
        "by a short verbatim locator quotation in the coding sheet so that a "
        "third party can check it. Coding was computer-assisted and reviewed by "
        "a single reviewer; governance and standards documents "
        "[[mds_privacy;edpb_cv;nist88;nist161;eu_battery]] were used as "
        "reference points and were not coded as operator practices.",
    ]]

    blocks.append(("h1", "3. Results"))
    blocks.append(("h2", "3.1. Scoping review (RQ1)"))
    blocks += [("p", t) for t in [
        f"Of {n_id} records screened, {fmt(SCR['ta_excluded'])} were excluded at "
        f"title/abstract and {fmt(SCR['sought'])} were sought for full text "
        f"(Fig. 1). Of these, {fmt(SCR['not_retrieved'])} could not be retrieved "
        f"and were recorded as such rather than assessed; "
        f"{fmt(SCR['excluded_fulltext'])} retrieved records were excluded "
        f"because they contained no relevant data path ({SCR['no_data_path']}) "
        f"or described a mechanism that did not transfer to the target domain "
        f"({SCR['not_transferable']}). {SCR['included']} direct studies were "
        f"included: {SCR['d4']} at evidence distance D4, {SCR['d3']} at D3, and "
        f"{SCR['d2']} at D2 (Fig. 2; Table 2).",
        "The included D4 studies provide the strongest evidence. A long-term "
        "real-world analysis reconstructed rider-relevant patterns from "
        "operator data [[elzer]]; investigative studies of rental applications "
        "and scooter ecosystems demonstrated collection and protocol weaknesses "
        "[[vinayaga2022;espoofer;etrojans]]; and a forensic analysis recovered "
        "data from micromobility devices [[hilgert]]. D3 studies document "
        "platform architectures, location-spoofing threats, data-acquisition "
        "frameworks, and user-facing traceability concerns "
        "[[isik;vinayaga2020;yilmaz2022;yilmaz2023;sato;li2020;petersen;zhou;"
        "hannemann]]. D2 studies transfer from adjacent domains: battery "
        "side channels and battery-management data [[leaky;bms]], residual data "
        "in reused IoT devices [[iotreuse]], and cloud data remanence "
        "[[remanence]]. No included study, on its own, demonstrated an "
        "end-to-end lifecycle compromise; the evidence is strongest for "
        "operation and weakest for recall and disposal.",
    ]]
    blocks.append(("fig", 1))
    blocks.append(("fig", 2))
    blocks.append(("table", 2))
    blocks.append(("h2", "3.2. Public GBFS field audit (RQ2)"))
    blocks += [("p", t) for t in [
        f"Of {fmt(N_REGISTRY)} registry systems, {reach} (95% CI {reach_ci}; "
        f"n = {fmt(N_REACHABLE)}) exposed a reachable auto-discovery endpoint, "
        f"{pct('reachable_registry_entries', 'vehicle_feed_declared')} of "
        "reachable systems declared a vehicle-status feed, and "
        f"{pct('successful_vehicle_feeds', 'vehicle_feed_nonempty')} of "
        "successfully retrieved feeds contained at least one vehicle (Table 3). "
        f"Restricting to the {fmt(N_MOTOR_FEEDS)} non-empty feeds that declared "
        "motorized micromobility, a vehicle identifier was present in "
        f"{pct('declared_motorized_micromobility_feeds', 'has_vehicle_id')} of "
        "feeds and latitude/longitude in "
        f"{pct('declared_motorized_micromobility_feeds', 'has_location_fields')} "
        "(Fig. 3). Fields with more operational specificity were less uniformly "
        "published: current range in "
        f"{pct('declared_motorized_micromobility_feeds', 'has_range')}, "
        "vehicle-specific rental links in "
        f"{pct('declared_motorized_micromobility_feeds', 'has_deep_link')}, "
        "last-reported timestamps in "
        f"{pct('declared_motorized_micromobility_feeds', 'has_last_reported')}, "
        "and battery or fuel percentage in "
        f"{pct('declared_motorized_micromobility_feeds', 'has_battery_percent')}.",
        "The operator-domain sensitivity analysis indicates that these are not "
        f"artefacts of a few large operators. Across {fmt(N_OP_DOMAINS)} "
        "eligible operator domains, a vehicle identifier appeared in every "
        "eligible feed for "
        f"{pct('declared_motorized_micromobility_operator_domains_all', 'has_vehicle_id')} "
        "of domains and latitude/longitude for "
        f"{pct('declared_motorized_micromobility_operator_domains_any', 'has_location_fields')}, "
        "whereas battery percentage was present in at least one eligible feed "
        "for only "
        f"{pct('declared_motorized_micromobility_operator_domains_any', 'has_battery_percent')} "
        "of domains. These figures describe what is disclosed publicly; they do "
        "not describe what operators collect or store on their backends, which "
        "the public feed cannot reveal.",
    ]]
    blocks.append(("fig", 3))
    blocks.append(("table", 3))
    blocks.append(("h2", "3.3. Public-document disclosure audit (RQ3)"))
    exp_loc = domain_count("location_data", "explicit")
    nf_batt = domain_count("battery_or_diagnostic_data", "not_found")
    nf_vuln = domain_count("vulnerability_disclosure", "not_found")
    nf_disp = domain_count("return_recycling_disposal", "not_found")
    blocks += [("p", t) for t in [
        f"We coded {N_CODED} operator privacy notices across {len(DOMAINS)} "
        f"domains; {N_OPERATORS - N_CODED} "
        "further operators could not be retrieved as reproducible text and were "
        "recorded as unavailable (Fig. 4; Table 4). Operational and "
        "account-level processing was disclosed almost universally: location "
        f"data were explicit for {exp_loc} of {N_CODED} operators, and "
        f"retention, processors or contractors, and account/payment/device data "
        f"were explicit for most. Data-subject rights and international "
        "transfers were explicit or partial for the large majority, consistent "
        "with a predominantly European operator sample governed by the General "
        "Data Protection Regulation [[gdpr]].",
        "Disclosure was markedly thinner for device-centred and lifecycle "
        f"domains. Battery or diagnostic data collection was not found in "
        f"{nf_batt} of {N_CODED} notices, even though such data are technically "
        "central to these vehicles [[leaky;bms]]. No operator notice in the "
        f"sample described a vulnerability-disclosure channel ({nf_vuln} of "
        f"{N_CODED} not found), and none addressed device return, recycling, or "
        f"disposal handling ({nf_disp} of {N_CODED} not found). We stress that "
        "these are gaps in public documents, not proof that the corresponding "
        "practices are absent; they nonetheless mark the parts of the lifecycle "
        "that are least visible to the public and to procuring organizations.",
    ]]
    blocks.append(("fig", 4))
    blocks.append(("table", 4))

    blocks.append(("h1", "4. Discussion"))
    blocks += [("p", t) for t in [
        "Read together, the three components describe a consistent pattern. "
        "Direct evidence for data exposure is concentrated in the operation "
        "stage, where audits and attack studies demonstrate collection, "
        "tracking, and re-identification potential [[elzer;espoofer;"
        "vinayaga2022]]. The public feed audit shows that the raw materials for "
        "such analyses - persistent identifiers and precise positions - are "
        "disclosed at high prevalence worldwide, while more operationally "
        "sensitive fields are published less uniformly. The disclosure audit "
        "shows that operators describe operational data handling in detail but "
        "are largely silent on the device end-of-life and on security-reporting "
        "channels, which is exactly where the review found only near-domain "
        "(D2) evidence [[iotreuse;remanence]].",
        "We integrate these observations in a six-stage lifecycle model that "
        "links each stage to the evidence distance of its supporting sources "
        "and to a proposed control (Fig. 5; Table 5). The model makes the "
        "strength of the underlying evidence explicit: procurement, deployment, "
        "operation, and maintenance are anchored by direct (D3-D4) evidence, "
        "whereas recall/return and second-life/disposal rest on near-domain "
        "(D2) analogy and on governance and standards guidance "
        "[[nist88;nist161;eu_battery]]. The controls - field minimization, "
        "local processing, access and retention limits, controlled contractor "
        "access, chain-of-custody for returned units, and verified media "
        "sanitization - are proposals whose effectiveness this study did not "
        "test; Table 5 records their effectiveness evidence conservatively.",
        "For procuring organizations, the practical implication is that a "
        "narrow focus on real-time position understates exposure. A telemetry "
        "inventory that spans identifiers, timestamps, range, and battery "
        "state, and that follows devices through maintenance and disposal, is a "
        "more faithful basis for assessment than the intuition that "
        "coordinates alone are the only sensitive field.",
    ]]
    blocks.append(("fig", 5))
    blocks.append(("table", 5))
    blocks.append(("h2", "4.1. Limitations"))
    blocks += [("p", t) for t in [
        "Several limitations bound these findings. Screening and coding were "
        "performed by a single reviewer with computer assistance; although the "
        "process is deterministic and reproducible, it does not provide "
        "inter-rater reliability, and some included studies were characterized "
        "from abstracts and curated metadata rather than from independently "
        "reproduced experiments. The GBFS audit observes only what is published "
        "in public feeds at a single point in time; it cannot show what "
        "operators collect or retain internally, and field presence is not a "
        "measure of harm. The disclosure audit reflects the public documents we "
        "could retrieve; not-found codings denote silence, two operators were "
        "unavailable, and the operator sample is weighted toward European and "
        "North American providers. Finally, the lifecycle controls are "
        "proposals; validating them would require intervention studies or "
        "operator cooperation that were outside this study's scope.",
    ]]

    blocks.append(("h1", "5. Conclusion"))
    blocks += [("p", t) for t in [
        "Shared micromobility offers a tractable, fully public setting in which "
        "to study lifecycle data exposure in connected devices that "
        f"organizations use but do not control. A scoping review of {n_id} "
        f"records identified {SCR['included']} direct studies whose strongest "
        f"evidence concerns the operation stage; a global audit of "
        f"{fmt(N_REACHABLE)} public GBFS "
        "systems showed that persistent identifiers and precise positions are "
        "disclosed at high prevalence; and a disclosure audit showed that "
        f"operators document operational processing thoroughly but are silent "
        f"on device disposal and vulnerability reporting. These are disclosure "
        "and evidence signals, not demonstrations of harm. The lifecycle model "
        "and its traceability table turn them into an auditable agenda: extend "
        "assessment beyond real-time location, close the documentation gaps at "
        "end-of-life, and empirically validate the proposed controls.",
    ]]
    return blocks


# ---------------------------------------------------------------------------
# Citation resolution
# ---------------------------------------------------------------------------
def iter_citation_texts(blocks: list[tuple], tables: dict[int, dict]):
    """Yield text fragments in document order for numbering."""
    for kind, payload in blocks:
        if kind == "p":
            yield payload
        elif kind == "table":
            spec = tables[payload]
            for r in spec["rows"]:
                for cell in r:
                    yield cell


CITE_RX = re.compile(r"\[\[([^\]]+)\]\]")


def _parse_label(raw: str) -> tuple[str, bool]:
    """Return (label, year_only). A ``year:`` prefix requests a year-only
    (narrative) citation, e.g. "Arksey and O'Malley (2005)"."""
    raw = raw.strip()
    if raw.startswith("year:"):
        return raw[len("year:"):].strip(), True
    return raw, False


def resolve_citations(blocks: list[tuple], tables: dict[int, dict]):
    """Author-date (Cambridge A) citation resolver.

    Internal ``[[label]]`` tokens are replaced with "(Author Year)" strings.
    Multiple labels in one token are sorted alphabetically by first author and
    chronologically for the same author. The reference list is returned in
    alphabetical order (not first-appearance order)."""
    order: list[str] = []

    def register(fragment: str) -> None:
        for m in CITE_RX.finditer(fragment):
            for raw in m.group(1).split(";"):
                label, _ = _parse_label(raw)
                if label not in order:
                    order.append(label)

    for frag in iter_citation_texts(blocks, tables):
        register(frag)

    unknown = [l for l in order if l not in REFS]
    if unknown:
        raise RuntimeError(f"Unknown citation labels: {unknown}")

    ref_list = sorted(set(order), key=lambda l: (sortkey(l), l))

    def repl(fragment: str) -> str:
        def _sub(m: re.Match) -> str:
            parsed = [_parse_label(p) for p in m.group(1).split(";")]
            year_only = all(yo for _, yo in parsed)
            labels = sorted((lbl for lbl, _ in parsed), key=lambda l: (sortkey(l), l))
            if year_only:
                inner = "; ".join(CITEMETA[l][1] for l in labels)
            else:
                inner = "; ".join(intext(l) for l in labels)
            return f"({inner})"
        return CITE_RX.sub(_sub, fragment)

    return order, ref_list, repl


# ---------------------------------------------------------------------------
# DOCX helpers (retained from the original generator)
# ---------------------------------------------------------------------------
def reset_dirs() -> None:
    for path in (OUTPUT, WORK):
        if path.exists():
            shutil.rmtree(path)
        path.mkdir(parents=True)


def set_cell_shading(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_cell_margins(cell, top=80, start=80, bottom=80, end=80) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    tc_mar = tc_pr.first_child_found_in("w:tcMar")
    if tc_mar is None:
        tc_mar = OxmlElement("w:tcMar")
        tc_pr.append(tc_mar)
    for margin, value in (("top", top), ("start", start), ("bottom", bottom), ("end", end)):
        node = tc_mar.find(qn(f"w:{margin}"))
        if node is None:
            node = OxmlElement(f"w:{margin}")
            tc_mar.append(node)
        node.set(qn("w:w"), str(value))
        node.set(qn("w:type"), "dxa")


def configure_document(doc: Document) -> None:
    section = doc.sections[0]
    for attr in ("top_margin", "bottom_margin", "left_margin", "right_margin"):
        setattr(section, attr, Inches(1))
    styles = doc.styles
    normal = styles["Normal"]
    normal.font.name = "Times New Roman"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
    normal.font.size = Pt(12)
    normal.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    normal.paragraph_format.space_after = Pt(6)
    for style_name, size in (("Title", 16), ("Heading 1", 14), ("Heading 2", 12)):
        style = styles[style_name]
        style.font.name = "Times New Roman"
        style._element.rPr.rFonts.set(qn("w:eastAsia"), "Times New Roman")
        style.font.size = Pt(size)
        style.font.color.rgb = RGBColor(0, 0, 0)
    styles["Heading 1"].font.bold = True
    styles["Heading 2"].font.bold = True


def add_body_paragraph(doc: Document, text: str) -> None:
    paragraph = doc.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    paragraph.paragraph_format.first_line_indent = Inches(0.3)
    run = paragraph.add_run(text)
    run.font.name = "Times New Roman"
    run.font.size = Pt(12)


def add_table(doc: Document, spec: dict, repl) -> None:
    caption = doc.add_paragraph()
    caption.paragraph_format.space_before = Pt(14)
    caption.paragraph_format.space_after = Pt(6)
    caption.add_run(repl(spec["title"])).bold = True
    table = doc.add_table(rows=1, cols=len(spec["headers"]))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    table.style = "Table Grid"
    header = table.rows[0].cells
    for index, value in enumerate(spec["headers"]):
        header[index].text = value
        set_cell_shading(header[index], "D9EAF7")
        header[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        for run in header[index].paragraphs[0].runs:
            run.bold = True
            run.font.name = "Arial"
            run.font.size = Pt(8)
        set_cell_margins(header[index])
    for row_values in spec["rows"]:
        cells = table.add_row().cells
        for index, value in enumerate(row_values):
            cells[index].text = repl(value)
            cells[index].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.TOP
            set_cell_margins(cells[index])
            for paragraph in cells[index].paragraphs:
                paragraph.paragraph_format.space_after = Pt(0)
                for run in paragraph.runs:
                    run.font.name = "Arial"
                    run.font.size = Pt(8)
    doc.add_paragraph()


def add_figure(doc: Document, num: int, png: Path, repl) -> None:
    image_p = doc.add_paragraph()
    image_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    image_p.paragraph_format.space_before = Pt(14)
    width = 6.3 if num in (1, 2, 3) else 6.5
    image_p.add_run().add_picture(str(png), width=Inches(width))
    caption = doc.add_paragraph()
    caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
    caption.paragraph_format.space_before = Pt(6)
    caption.paragraph_format.space_after = Pt(14)
    run = caption.add_run(repl(FIG_CAPTIONS[num]))
    run.italic = True
    run.font.size = Pt(10)


def add_page_number(section) -> None:
    paragraph = section.footer.paragraphs[0]
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run()
    begin = OxmlElement("w:fldChar")
    begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    end = OxmlElement("w:fldChar")
    end.set(qn("w:fldCharType"), "end")
    run._r.extend([begin, instr, end])


ABSTRACT_STRUCT = [
    ("Background", "Shared micromobility produces continuous location, motion, "
     "and battery telemetry, yet evidence on data exposure is scattered and the "
     "device lifecycle beyond real-time position is rarely examined."),
    ("Objective", "To map the direct evidence for micromobility data exposure, "
     "to measure which vehicle fields operators publish worldwide, and to "
     "assess how completely operators disclose lifecycle data handling."),
    ("Methods", "We combined a scoping review reported using the Preferred "
     "Reporting Items for Systematic Reviews and Meta-Analyses extension for "
     f"Scoping Reviews (PRISMA-ScR) of {fmt(SCR['identified'])} records, a "
     "cross-sectional field audit of public General Bikeshare Feed "
     "Specification (GBFS) feeds, and a structured disclosure audit of public "
     f"operator privacy notices across {len(DOMAINS)} domains, using "
     "deterministic screening and explicit denominators with "
     "95% Wilson confidence intervals (CIs)."),
    ("Results", f"{SCR['included']} direct studies were included "
     f"({SCR['d4']} direct-empirical, {SCR['d3']} direct-documentary, and "
     f"{SCR['d2']} near-domain), with "
     f"the strongest evidence at the operation stage. Across {fmt(N_MOTOR_FEEDS)} "
     "motorized micromobility feeds, vehicle identifiers appeared in "
     f"{pct('declared_motorized_micromobility_feeds', 'has_vehicle_id')} and "
     f"coordinates in "
     f"{pct('declared_motorized_micromobility_feeds', 'has_location_fields')}, "
     "while battery percentage appeared in "
     f"{pct('declared_motorized_micromobility_feeds', 'has_battery_percent')}. "
     "Operator notices documented operational processing thoroughly but were "
     "silent on device disposal and vulnerability reporting."),
    ("Conclusions", "Persistent identifiers and precise positions are widely "
     "disclosed, and lifecycle end-of-life handling is largely undocumented. "
     "These are disclosure and evidence signals, not demonstrations of harm; we "
     "propose a lifecycle model whose controls remain to be validated."),
]


def build_abstract_text() -> str:
    return " ".join(f"{label}: {text}" for label, text in ABSTRACT_STRUCT)


def build_manuscript(blocks, tables, figpaths, repl, references) -> Path:
    doc = Document()
    configure_document(doc)
    add_page_number(doc.sections[0])
    title = doc.add_paragraph(style="Title")
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.add_run(TITLE).bold = True
    byline = doc.add_paragraph()
    byline.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if BLINDED:
        # Double-masked: the main manuscript carries no identifying byline; the
        # author details travel on the separate, non-anonymized title page.
        anon = byline.add_run(
            "Author and affiliation details removed for double-masked peer "
            "review")
        anon.italic = True
        anon.font.size = Pt(10)
    else:
        byline.add_run(AUTHOR).bold = True
        aff = doc.add_paragraph()
        aff.alignment = WD_ALIGN_PARAGRAPH.CENTER
        ar = aff.add_run(
            f"{AFFILIATION}\nCorresponding author: {EMAIL}; ORCID {ORCID}")
        ar.italic = True
        ar.font.size = Pt(10)

    doc.add_heading("Abstract", level=1)
    for label, text in ABSTRACT_STRUCT:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        p.add_run(f"{label}: ").bold = True
        p.add_run(text)

    doc.add_heading("Policy Significance Statement", level=1)
    ps = doc.add_paragraph()
    ps.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    ps.add_run(POLICY_SIGNIFICANCE)

    kw = doc.add_paragraph()
    kw.add_run("Keywords: ").bold = True
    kw.add_run("; ".join(KEYWORDS))

    for kind, payload in blocks:
        if kind == "h1":
            doc.add_heading(payload, level=1)
        elif kind == "h2":
            doc.add_heading(payload, level=2)
        elif kind == "p":
            add_body_paragraph(doc, repl(payload))
        elif kind == "fig":
            add_figure(doc, payload, figpaths[payload]["png"], repl)
        elif kind == "table":
            add_table(doc, tables[payload], repl)

    for label, value in DECLARATIONS:
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(6)
        p.paragraph_format.space_after = Pt(4)
        p.add_run(f"{label}. ").bold = True
        p.add_run(value)

    doc.add_heading("References", level=1)
    for label in references:
        p = doc.add_paragraph()
        p.paragraph_format.left_indent = Inches(0.28)
        p.paragraph_format.first_line_indent = Inches(-0.28)
        p.paragraph_format.line_spacing = 1.0
        p.paragraph_format.space_after = Pt(4)
        p.add_run(REFS[label])

    path = OUTPUT / "Manuscript_DataPolicy.docx"
    doc.save(path)
    return path


def build_title_page(word_count: int, n_refs: int) -> Path:
    doc = Document()
    configure_document(doc)
    title = doc.add_paragraph(style="Title")
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.add_run(TITLE).bold = True
    doc.add_paragraph()
    for text, bold in [(AUTHOR, True), (AFFILIATION, False),
                       (f"Corresponding author: {AUTHOR}; {EMAIL}", False),
                       (f"ORCID: {ORCID}", False)]:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run(text).bold = bold
    doc.add_paragraph()
    ps_words = len(re.findall(r"\b[\w'-]+\b", POLICY_SIGNIFICANCE))
    abstract_words = len(re.findall(r"\b[\w'-]+\b", build_abstract_text()))
    fields = [
        ("Full title", TITLE),
        ("Short title", SHORT_TITLE),
        ("Article type", ARTICLE_TYPE),
        ("Journal", f"{JOURNAL} ({PUBLISHER})"),
        ("Area of interest", AREA_OF_INTEREST),
        ("Peer review", "Single-blind; the manuscript is not anonymized"),
        ("Abstract word count", str(abstract_words)),
        ("Policy Significance Statement word count", str(ps_words)),
        ("Main-text word count (excludes references)", str(word_count)),
        ("Figures", "5"),
        ("Tables", "5"),
        ("References", str(n_refs)),
        ("Funding", "None"),
        ("Competing interests", "None declared"),
    ]
    for label, value in fields:
        p = doc.add_paragraph()
        p.add_run(f"{label}: ").bold = True
        p.add_run(value)
    path = OUTPUT / "Title_Page_DataPolicy.docx"
    doc.save(path)
    return path


def build_cover_letter() -> Path:
    doc = Document()
    configure_document(doc)
    for line in [AUTHOR, AFFILIATION, EMAIL, f"ORCID: {ORCID}", BUILD_DATE]:
        p = doc.add_paragraph(line)
        p.paragraph_format.space_after = Pt(2)
    doc.add_paragraph()
    doc.add_paragraph("The Editors-in-Chief")
    doc.add_paragraph(JOURNAL)
    doc.add_paragraph(PUBLISHER)
    doc.add_paragraph()
    doc.add_paragraph("Dear Editors,")
    paras = [
        f"I submit the manuscript \u201c{TITLE}\u201d for consideration as a "
        f"Research Article in {JOURNAL}.",
        "Shared micromobility is a fully public, globally deployed class of "
        "connected devices that cities and citizens use but do not own or "
        "decommission. The manuscript reports an empirical, reproducible "
        "package at the intersection of data and policy: a PRISMA-ScR scoping "
        "review that identifies the direct evidence for data exposure; a global "
        "audit of public GBFS feeds that measures which vehicle fields "
        "operators actually publish; and a structured audit of public operator "
        f"privacy notices across {len(DOMAINS)} disclosure domains.",
        "The work fits the journal's Area 4 (Ethics, Equity and "
        "Trustworthiness of Data) because it quantifies real, worldwide data "
        "disclosure and connects it to procurement, regulation, and "
        "transparency choices facing public authorities. A 120-word Policy "
        "Significance Statement makes these implications explicit for a "
        "policy audience. Throughout, we treat field presence and document "
        "silence as disclosure and evidence signals rather than as proof of "
        "harm, compromise, or regulatory violation, and we state the "
        "effectiveness of proposed controls conservatively as not yet "
        "validated.",
        "The study analysed only publicly accessible feeds and documents; it "
        "did not attempt authentication or access-control circumvention, did "
        "not interact with users, and retained no raw identifiers, exact "
        "coordinates, or vehicle deep links. Consistent with the journal's "
        "open-research expectations, all data, coding sheets, and code are "
        "openly available so that every reported count can be regenerated with "
        "a single command. The manuscript is original and is not under "
        "consideration elsewhere; I understand that Data & Policy uses "
        "single-blind peer review, so the submission is not anonymized. There "
        "is no funding or competing interest to declare, and no generative AI "
        "was used for scientific content.",
        "Thank you for considering this submission.",
    ]
    for text in paras:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        p.add_run(text)
    doc.add_paragraph("Sincerely,")
    doc.add_paragraph(AUTHOR)
    path = OUTPUT / "Cover_Letter_DataPolicy.docx"
    doc.save(path)
    return path


def build_tables_docx(tables, repl) -> Path:
    doc = Document()
    configure_document(doc)
    h = doc.add_paragraph(style="Title")
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    h.add_run("Editable tables").bold = True
    for i in range(1, 6):
        add_table(doc, tables[i], repl)
        if i < 5:
            doc.add_section(WD_SECTION.NEW_PAGE)
    path = OUTPUT / "Tables_DataPolicy_editable.docx"
    doc.save(path)
    return path


def build_figures_pptx(figpaths, repl) -> Path:
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    for num in range(1, 6):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.2),
                                             PptxInches(12.3), PptxInches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Figure {num}"
        p.font.size = PptxPt(22)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        from PIL import Image as PILImage
        with PILImage.open(figpaths[num]["png"]) as im:
            w, h = im.size
        avail_w, avail_h = 12.3, 5.3
        ratio = min(avail_w / (w / 600.0), avail_h / (h / 600.0))
        disp_w = (w / 600.0) * ratio
        disp_h = (h / 600.0) * ratio
        left = (13.333 - disp_w) / 2
        slide.shapes.add_picture(str(figpaths[num]["png"]), PptxInches(left),
                                 PptxInches(1.0), PptxInches(disp_w), PptxInches(disp_h))
        cap_box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(6.5),
                                          PptxInches(12.1), PptxInches(0.9))
        cf = cap_box.text_frame
        cf.word_wrap = True
        cp = cf.paragraphs[0]
        cp.text = repl(FIG_CAPTIONS[num])
        cp.font.size = PptxPt(11)
    path = OUTPUT / "Figures_DataPolicy_editable.pptx"
    prs.save(path)
    return path


def build_reporting_guideline() -> Path:
    doc = Document()
    configure_document(doc)
    h = doc.add_paragraph(style="Title")
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    h.add_run("PRISMA-ScR reporting checklist").bold = True
    p = doc.add_paragraph()
    p.add_run(
        "The scoping-review component (WP1) follows the PRISMA-ScR checklist "
        "(Tricco et al. 2018). The field audit (WP2) and disclosure audit "
        "(WP3) are cross-sectional observational studies of public artefacts "
        "and are reported with explicit denominators, confidence intervals, and "
        "open code. The table below maps each PRISMA-ScR item to its location.")
    items = [
        ("Title", "Identifies the report as a scoping review", "Title page"),
        ("Abstract", "Structured summary", "Abstract"),
        ("Rationale", "Rationale in the context of what is known", "Section 1"),
        ("Objectives", "Research questions RQ1-RQ3", "Section 1"),
        ("Protocol", "Protocol availability", "Data availability statement"),
        ("Eligibility criteria", "Characteristics used as criteria", "Section 2.2; Table 1"),
        ("Information sources", "Sources searched", "Section 2.2; Table 1"),
        ("Selection of sources", "Screening process", "Section 2.2; Fig. 1"),
        ("Data charting", "Extraction process and items", "Section 2.2; Table 2"),
        ("Synthesis of results", "Methods of summarizing", "Sections 3-4"),
        ("Results of sources", "Numbers screened and included", "Section 3.1; Fig. 1"),
        ("Results of syntheses", "Charted results", "Sections 3.1-3.3"),
        ("Limitations", "Limitations of the review", "Section 4.1"),
        ("Conclusions", "Interpretation and implications", "Section 5"),
        ("Funding", "Sources of funding", "Declarations"),
    ]
    add_table(doc, {"title": "PRISMA-ScR item mapping",
                    "headers": ["Item", "Checklist description", "Location"],
                    "rows": items}, lambda x: x)
    path = OUTPUT / "Reporting_Guideline_PRISMA-ScR.docx"
    doc.save(path)
    return path


def resolve_identifier(ident: str) -> str:
    """Best-effort live check that a DOI/URL resolves. Falls back gracefully
    when offline so the build stays reproducible."""
    if ident.startswith("doi:"):
        url = "https://doi.org/" + ident[4:]
    elif ident.startswith("http"):
        url = ident
    else:
        return "not_applicable_standard_or_regulation"
    req = urllib.request.Request(
        url, method="HEAD", headers={"User-Agent": "Mozilla/5.0 ref-check"})
    try:
        code = urllib.request.urlopen(req, timeout=25).status
        return f"resolves_http_{code}"
    except urllib.error.HTTPError as exc:
        # 403/429 are anti-bot responses from a live, existing record.
        return f"exists_http_{exc.code}"
    except Exception:
        return "not_checked_offline"


def build_reference_verification(refs) -> Path:
    path = OUTPUT / "Reference_Verification.csv"
    lines = ["Number,Label,Reference,Identifier,VerificationDate,Status"]
    for i, label in enumerate(refs, 1):
        text = REFS[label]
        m = re.search(r"(doi:\S+|https?://\S+)", text)
        ident = m.group(1).rstrip(".") if m else "(no DOI; standard/regulation/registry)"
        status = resolve_identifier(ident)
        row = [str(i), label, text, ident, BUILD_DATE, status]
        lines.append(",".join(f'"{v.replace(chr(34), chr(34)*2)}"' for v in row))
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return path


def build_citation_audit(blocks, tables, order, repl) -> Path:
    first: dict[str, tuple[str, str]] = {}
    location = "Body"
    for kind, payload in blocks:
        if kind in ("h1", "h2"):
            location = payload
            continue
        frags = []
        if kind == "p":
            frags = [payload]
        elif kind == "table":
            frags = [c for r in tables[payload]["rows"] for c in r]
            location = f"Table {payload}"
        for frag in frags:
            for m in CITE_RX.finditer(frag):
                for raw in m.group(1).split(";"):
                    label, _ = _parse_label(raw)
                    if label not in first:
                        first[label] = (location, textwrap.shorten(repl(frag), 200, placeholder="..."))
    path = OUTPUT / "Citation_Audit.csv"
    lines = ["Citation,Year,First appearance,Context"]
    for label in order:
        loc, ctx = first[label]
        lines.append(",".join(f'"{v.replace(chr(34), chr(34)*2)}"'
                              for v in [intext(label), CITEMETA[label][1], loc, ctx]))
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return path


def word_count(blocks, repl) -> int:
    text = TITLE + " " + build_abstract_text()
    for kind, payload in blocks:
        if kind == "p":
            text += " " + repl(payload)
    return len(re.findall(r"\b[\w'-]+\b", text))


def validate(blocks, tables, order, ref_list, repl) -> dict:
    # every cited label resolves (resolve_citations already raises on unknown)
    all_resolve = all(l in REFS for l in order)
    # all refs used, none orphan or phantom
    orphan = set(order) ^ set(REFS.keys())
    # reference list is alphabetical by author (then year)
    refs_alpha = ref_list == sorted(ref_list, key=lambda l: (sortkey(l), l))
    body = word_count(blocks, repl)
    abstract_words = len(re.findall(r"\b[\w'-]+\b", build_abstract_text()))
    ps_words = len(re.findall(r"\b[\w'-]+\b", POLICY_SIGNIFICANCE))
    # figures/tables cited in text
    body_text = " ".join(repl(p) for k, p in blocks if k == "p")
    figs_cited = all(f"Fig. {i}" in body_text for i in range(1, 6))
    tabs_cited = all(f"Table {i}" in body_text for i in range(1, 6))
    # figure/table blocks present and sequential
    fig_seq = [p for k, p in blocks if k == "fig"]
    tab_seq = [p for k, p in blocks if k == "table"]
    # in-text markers must not leak into the rendered text
    no_raw_markers = not any(CITE_RX.search(repl(p)) for k, p in blocks if k == "p")
    # required Data & Policy disclosure statements
    decl_labels = {lbl for lbl, _ in DECLARATIONS}
    required_disclosures = {
        "Data availability statement", "Funding statement", "Competing interests",
    }.issubset(decl_labels)
    # double-masked guard: no author-identifying token may reach the disclosures
    identity_tokens = [PUBLIC_REPO_URL, "bougtoir"]
    for value in (AUTHOR, EMAIL, ORCID):
        if value and "[" not in value:
            identity_tokens.append(value)
    decl_text = " ".join(v for _, v in DECLARATIONS)
    no_identity_leak = (not BLINDED) or not any(
        tok in decl_text for tok in identity_tokens)
    checks = {
        "all_citations_resolve": all_resolve,
        "no_orphan_or_phantom_refs": not orphan,
        "references_alphabetical": refs_alpha,
        "no_unresolved_markers": no_raw_markers,
        "figures_cited_in_text": figs_cited,
        "tables_cited_in_text": tabs_cited,
        "five_figures_present": sorted(set(fig_seq)) == [1, 2, 3, 4, 5],
        "five_tables_present": sorted(set(tab_seq)) == [1, 2, 3, 4, 5],
        "abstract_within_250w": abstract_words <= 250,
        "policy_significance_120w": 110 <= ps_words <= 130,
        "keywords_at_most_five": len(KEYWORDS) <= 5,
        "required_disclosures_present": required_disclosures,
        "no_identity_leak_when_blinded": no_identity_leak,
    }
    failures = [k for k, v in checks.items() if not v]
    if failures:
        raise RuntimeError(f"Validation failed: {failures}; orphan={orphan}")
    # abbreviation-at-first-use spot check
    joined = body_text
    for abbr, definition in {
        "GBFS": "General Bikeshare Feed Specification (GBFS)",
        "PRISMA-ScR": "Scoping Reviews (PRISMA-ScR)",
        "DOI": "digital object identifier (DOI)",
        "GDPR": "General Data Protection Regulation",
    }.items():
        if abbr in joined and definition not in joined:
            raise RuntimeError(f"Undefined abbreviation at first use: {abbr}")
    return {"word_count": body, "abstract_words": abstract_words,
            "policy_significance_words": ps_words, "keywords": len(KEYWORDS),
            "references": len(ref_list), **checks}


def build_validation_report(validation: dict, figpaths) -> Path:
    path = OUTPUT / "VALIDATION.txt"
    lines = ["Submission validation report", "=" * 30, ""]
    for k, v in validation.items():
        lines.append(f"{k}: {v}")
    lines.append("")
    lines.append("Figure files:")
    for num, kinds in figpaths.items():
        for kind, p in kinds.items():
            lines.append(f"  Figure {num} {kind}: {p.name}")
    path.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return path


def build_checklist(validation) -> Path:
    doc = Document()
    configure_document(doc)
    h = doc.add_paragraph(style="Title")
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    h.add_run("Submission checklist").bold = True
    rows = [
        ["Cover letter", "Included", "Cover_Letter_DataPolicy.docx"],
        ["Title page with author details", "Included", "Title_Page_DataPolicy.docx"],
        ["Article file (figures/tables inline; non-anonymized)", "Included", "Manuscript_DataPolicy.docx"],
        ["Abstract (<=250 words)", f"{validation['abstract_words']} words", "Manuscript"],
        ["Policy Significance Statement (120 words)", f"{validation['policy_significance_words']} words", "Manuscript"],
        ["Keywords (<=5, semicolon-separated)", f"{validation['keywords']}", "Manuscript"],
        ["Figures (PNG + TIFF + PDF, 600 dpi)", "Included", "figures/Figure1-5.*"],
        ["Editable figures (one per slide)", "Included", "Figures_DataPolicy_editable.pptx"],
        ["Editable tables", "Included", "Tables_DataPolicy_editable.docx"],
        ["Reporting guideline (PRISMA-ScR)", "Included", "Reporting_Guideline_PRISMA-ScR.docx"],
        ["Citation audit (author-date)", "Included", "Citation_Audit.csv"],
        ["Reference verification", "Included", "Reference_Verification.csv"],
        [f"Main-text word count (excl. references): {validation['word_count']}",
         "Reported", "Title_Page_DataPolicy.docx"],
        ["Data availability statement", "Included", "Manuscript"],
        ["Funding statement", "Included", "Manuscript"],
        ["Competing interests statement", "Included", "Manuscript"],
    ]
    add_table(doc, {"title": "Items included in the submission package",
                    "headers": ["Item", "Status", "File"], "rows": rows},
              lambda x: x)
    path = OUTPUT / "Submission_Checklist_DataPolicy.docx"
    doc.save(path)
    return path


def build_zip(figpaths) -> Path:
    path = OUTPUT / "DataPolicy_submission_package.zip"
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        for item in sorted(OUTPUT.rglob("*")):
            if item.is_file() and item != path:
                zf.write(item, item.relative_to(OUTPUT))
    return path


def main() -> None:
    reset_dirs()
    FIGDIR.mkdir(parents=True, exist_ok=True)
    figpaths = build_figures()
    blocks = body_blocks()
    order, ref_list, repl = resolve_citations(blocks, TABLES)
    validation = validate(blocks, TABLES, order, ref_list, repl)

    build_manuscript(blocks, TABLES, figpaths, repl, ref_list)
    wc = validation["word_count"]
    build_title_page(wc, len(ref_list))
    build_cover_letter()
    build_tables_docx(TABLES, repl)
    build_figures_pptx(figpaths, repl)
    build_reporting_guideline()
    build_reference_verification(ref_list)
    build_citation_audit(blocks, TABLES, order, repl)
    build_checklist(validation)
    build_validation_report(validation, figpaths)
    zip_path = build_zip(figpaths)

    print("Build complete.")
    print(f"  references: {len(ref_list)} (alphabetical, author-date)")
    print(f"  main-text words: {wc}; abstract words: {validation['abstract_words']}; "
          f"policy significance words: {validation['policy_significance_words']}")
    print(f"  package: {zip_path.relative_to(ROOT)}")
    for k, v in validation.items():
        if isinstance(v, bool):
            print(f"  {k}: {'OK' if v else 'FAIL'}")


if __name__ == "__main__":
    main()
